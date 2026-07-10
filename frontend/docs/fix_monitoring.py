# -*- coding: utf-8 -*-
"""렌더러가 누락시키는 손상된 '실시간 상담 모니터링' 원본 슬라이드를
깨끗한 버전(빈 스크린샷 + 정의/핵심기능/POINT)으로 교체. 0616 덮어쓰기."""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu

FILE = "GenOnLIFE_AICC_데모_0616.pptx"
TARGET_TITLE = "실시간 상담 모니터링"

NAVY = RGBColor(0x0F, 0x34, 0x68); BLUE = RGBColor(0x2F, 0x6B, 0xB0)
MINT = RGBColor(0x15, 0xC2, 0xA2); ADMIN = RGBColor(0x0E, 0x9B, 0x86)
LIGHT = RGBColor(0xF2, 0xF8, 0xFF); LIGHTER = RGBColor(0xF7, 0xFA, 0xFE)
INK = RGBColor(0x10, 0x23, 0x3F); GRAY = RGBColor(0x5B, 0x6B, 0x80)
FAINT = RGBColor(0xA8, 0xB6, 0xC8); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCD, 0xD9, 0xE8); LTBLUE = RGBColor(0x9F, 0xC4, 0xEE)
FONT = "Apple SD Gothic Neo"

prs = Presentation(FILE)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(r, s, c, b=False, i=False):
    r.font.size = Pt(s); r.font.color.rgb = c; r.font.bold = b; r.font.italic = i; r.font.name = FONT


def rect(sl, l, t, w, h, fill, line=None, round_=False):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                              Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(sl, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, ital) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, ital)
    return tb


def role_chip(sl, l, t, text, col):
    w = 0.28 + 0.135 * len(text)
    shp = rect(sl, l, t, w, 0.4, WHITE, line=col, round_=True)
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set(r, 13, col, True)


def clean_slide(title, role, route, definition, features, point):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW.inches, SH.inches, WHITE)
    rect(s, 0, 0, SW.inches, 1.32, NAVY)
    rect(s, 0, 1.32, SW.inches, 0.05, MINT)
    textbox(s, 0.62, 0.26, 9.6, 0.7, [[(title, 26, WHITE, True, False)]])
    textbox(s, 0.64, 0.92, 9.6, 0.34, [[(route, 11.5, LTBLUE, False, False)]])
    role_chip(s, 11.55, 0.42, role + " 화면", ADMIN if role == "관리자" else NAVY)
    textbox(s, 12.2, 6.95, 0.9, 0.4, [[("00 / 00", 10, GRAY, False, False)]], align=PP_ALIGN.RIGHT)
    ph = rect(s, 0.4, 1.55, 7.5, 4.5, LIGHTER, line=BORDER, round_=True)
    tf = ph.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "화면 스크린샷 영역"; _set(r, 13, FAINT, True)
    RX, RE = 8.1, 12.9; rw = RE - RX
    textbox(s, RX, 1.6, 2, 0.34, [[("정의", 12, MINT, True, False)]])
    textbox(s, RX, 1.96, rw, 0.9, [[(definition, 13, GRAY, False, False)]])
    textbox(s, RX, 2.86, 4, 0.34, [[("핵심 기능", 12, BLUE, True, False)]])
    fy = 3.28
    for head, body in features:
        rect(s, RX + 0.04, fy + 0.06, 0.12, 0.12, BLUE)
        textbox(s, RX + 0.33, fy - 0.04, rw - 0.33, 0.72,
                [[(head + "  ", 13, INK, True, False), (body, 11, GRAY, False, False)]])
        fy += 0.72
    box = rect(s, 0.62, 6.05, 12.1, 0.82, LIGHT, line=RGBColor(0xBA, 0xD6, 0xF4), round_=True)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "POINT  "; _set(r, 12, MINT, True)
    r = p.add_run(); r.text = point; _set(r, 13, NAVY, False)
    return s


NEW = ("실시간 상담 모니터링 (관리자)", "관리자", "/realtime-monitoring",
       "실시간 코칭을 위한 관할 상담사 상담 모니터링 화면입니다.",
       [("관할 상담사 선택", "상담 중 상담사 선택 → AI 위험·주의 신호 표시"),
        ("실시간 STT 청취", "선택 상담사의 실시간 상담 내용을 읽기 전용으로 모니터링"),
        ("고객 정보·RAG", "고객 업무정보와 약관·업무 기준을 함께 확인"),
        ("실시간 코칭 채팅", "상담사에게 즉시 코칭 전달(귓속말/전체 안내)")],
       "관리자가 상담 현장을 실시간으로 보며 그 자리에서 코칭하는 메인 작업 화면입니다.")


def slide_title(slide):
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t and Emu(sh.top or 0).inches < 0.9 and "/" not in t and len(t) > 3:
                return t.split("\n")[0]
    return ""


def main():
    # 손상된 실시간 모니터링 슬라이드 위치 찾기
    lst = prs.slides._sldIdLst
    sldIds = list(lst)
    pos = None
    for i, (sldId, slide) in enumerate(zip(sldIds, prs.slides)):
        if slide_title(slide).startswith(TARGET_TITLE):
            pos = i; old_sldId = sldId; break
    if pos is None:
        print("대상 슬라이드 없음"); return

    # 새 깨끗한 슬라이드 생성 → 대상 위치로 이동, 손상본 제거
    clean_slide(*NEW)
    new_sldId = list(lst)[-1]
    lst.remove(new_sldId)
    lst.insert(pos, new_sldId)
    # 손상본 참조 제거(파트는 고아로 남지만 무해)
    rid = old_sldId.get(qn('r:id'))
    lst.remove(old_sldId)
    try: prs.part.drop_rel(rid)
    except Exception: pass

    # 페이지 번호 재부여
    slides = list(prs.slides); total = len(slides)
    pat = re.compile(r"^\s*\d{1,2}\s*/\s*\d{1,2}\s*$")
    for p, slide in enumerate(slides, 1):
        for sh in slide.shapes:
            if sh.has_text_frame and pat.match(sh.text_frame.text.strip()):
                sh.text_frame.paragraphs[0].runs[0].text = f"{p:02d} / {total:02d}"; break

    prs.save(FILE)
    print("교체 완료:", FILE, "| slides:", len(prs.slides._sldIdLst), "| pos:", pos + 1)


if __name__ == "__main__":
    main()
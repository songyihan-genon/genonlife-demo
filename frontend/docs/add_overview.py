# -*- coding: utf-8 -*-
"""슬라이드 2번에 '플랫폼 기능 요약(메뉴별 지원 업무)' 장표 추가. v2 → v3로 저장."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "GenOnLIFE_AICC_데모.pptx"
OUT = "GenOnLIFE_AICC_데모_v3.pptx"

NAVY = RGBColor(0x0F, 0x34, 0x68)
BLUE = RGBColor(0x2F, 0x6B, 0xB0)
MINT = RGBColor(0x15, 0xC2, 0xA2)
INK = RGBColor(0x10, 0x23, 0x3F)
GRAY = RGBColor(0x5B, 0x6B, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE6, 0xED, 0xF5)
FONT = "Apple SD Gothic Neo"

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height
DIV_GRAD = [(0, "0F3468"), (45, "2F8BFF"), (100, "15C2A2")]


def set_gradient(shape, stops, ang):
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = grad.makeelement(qn("a:gsLst"), {})
    for pos, rgb in stops:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": rgb}); gs.append(clr); gsLst.append(gs)
    grad.append(gsLst)
    grad.append(grad.makeelement(qn("a:lin"), {"ang": str(ang), "scaled": "1"}))
    ln = spPr.find(qn("a:ln"))
    (ln.addprevious(grad) if ln is not None else spPr.append(grad))


def rect(slide, l, t, w, h, fill, line=None, round_=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(slide, l, t, w, h, runs, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        for (txt, size, color, bold) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT
    return tb


def chip(slide, l, t, text, fill):
    w = 0.22 + 0.12 * len(text)
    shp = rect(slide, l, t, w, 0.36, fill, round_=True)
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    return w


def menu_card(slide, x, y, w, h, name, desc, accent):
    rect(slide, x, y, w, h, WHITE, line=BORDER, round_=True)
    rect(slide, x + 0.14, y + 0.2, 0.08, h - 0.4, accent, round_=True)  # 좌측 액센트 바
    textbox(slide, x + 0.34, y + 0.18, w - 0.5, 0.5, [[(name, 13.5, INK, True)]])
    textbox(slide, x + 0.34, y + 0.66, w - 0.55, h - 0.78, [[(desc, 10.5, GRAY, False)]])


# ── 새 슬라이드 생성 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
rect(slide, 0, 0, SW.inches, SH.inches, WHITE)
top = rect(slide, 0, 0, SW.inches, 0.22, MINT); set_gradient(top, DIV_GRAD, 0)
textbox(slide, 0.62, 0.5, 9.5, 0.7, [[("플랫폼 기능 요약", 28, INK, True)]])
textbox(slide, 0.64, 1.18, 11.5, 0.5, [[("메뉴별로 어떤 상담 업무를 지원하는가 — 상담사 · 관리자 관점", 14, GRAY, False)]])

# 상담사 영역
chip(slide, 0.62, 1.78, "상담사", NAVY)
agent_cards = [
    ("AI 상담 홈", "콜 인입 수신 · 고객정보·예상 문의 자동 로딩 · 개인 KPI · 업무 어시스턴트 검색"),
    ("실시간 고객 상담", "STT 받아쓰기 · 상담 가이드 · 지식검색(RAG) · 업무정보 연동 · 관리자 코칭"),
    ("상담 이력 · 후처리", "상담 이력 탐색 · 접촉이력 등록 · SMS 안내 발송 · 상담 검수 결과 열람"),
    ("고객 민원 탐지", "콜 STT 기반 불편 탐지 · 민원 발생 위험 평가 · VOC 자동 등록"),
]
n = len(agent_cards); AL, AR, GAP = 0.62, 12.71, 0.2
cw = (AR - AL - GAP * (n - 1)) / n
for k, (nm, ds) in enumerate(agent_cards):
    menu_card(slide, AL + k * (cw + GAP), 2.22, cw, 1.78, nm, ds, BLUE)

# 관리자 영역
chip(slide, 0.62, 4.2, "관리자", MINT)
admin_cards = [
    ("AI 상담 홈 (운영)", "실시간 운영 KPI · 상담사 상태·검수 분포·민원·품질 대시보드"),
    ("실시간 상담 모니터링", "관할 상담사 실시간 청취 · 고객정보·RAG 확인 · 실시간 코칭"),
    ("상담 검수", "AI 1차 검수 + 관리자 휴먼 2차 검수 · 코멘트·후속 조치"),
]
m = len(admin_cards); cw2 = (AR - AL - GAP * (m - 1)) / m
for k, (nm, ds) in enumerate(admin_cards):
    menu_card(slide, AL + k * (cw2 + GAP), 4.64, cw2, 1.78, nm, ds, MINT)

# 하단 한 줄 요약
band = rect(slide, 0.62, 6.66, 12.09, 0.56, RGBColor(0xF2, 0xF8, 0xFF), line=RGBColor(0xBA, 0xD6, 0xF4), round_=True)
btf = band.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE; btf.word_wrap = True
btf.margin_left = Inches(0.2)
bp = btf.paragraphs[0]
r = bp.add_run(); r.text = "콜 인입 → 실시간 상담 보조 → 후처리 자동화 → 검수·민원 관리까지, 상담 전 과정을 하나의 포털에서 지원"
r.font.size = Pt(12); r.font.color.rgb = NAVY; r.font.name = FONT

# ── 2번 위치로 이동 ──
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
sldIdLst.remove(ids[-1])
sldIdLst.insert(1, ids[-1])

prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
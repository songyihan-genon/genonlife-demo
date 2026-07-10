# -*- coding: utf-8 -*-
"""0616 덱에 '상담 품질 검수 (관리자)' 슬라이드 추가(관리자 관점).
add_new_screens.py 와 동일 레이아웃(텍스트 우측 + 빈 스크린샷 플레이스홀더).
실시간 상담 모니터링(13) 바로 뒤에 삽입. 페이지번호 재부여. 0616 덮어쓰기."""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FILE = "GenOnLIFE_AICC_데모_0616.pptx"

NAVY = RGBColor(0x0F, 0x34, 0x68)
BLUE = RGBColor(0x2F, 0x6B, 0xB0)
MINT = RGBColor(0x15, 0xC2, 0xA2)
ADMIN = RGBColor(0x0E, 0x9B, 0x86)
LIGHT = RGBColor(0xF2, 0xF8, 0xFF)
LIGHTER = RGBColor(0xF7, 0xFA, 0xFE)
INK = RGBColor(0x10, 0x23, 0x3F)
GRAY = RGBColor(0x5B, 0x6B, 0x80)
FAINT = RGBColor(0xA8, 0xB6, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCD, 0xD9, 0xE8)
LTBLUE = RGBColor(0x9F, 0xC4, 0xEE)
FONT = "Apple SD Gothic Neo"

prs = Presentation(FILE)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = FONT


def rect(slide, l, t, w, h, fill, line=None, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, italic) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, italic)
    return tb


def role_chip_outline(slide, l, t, text, col):
    w = 0.28 + 0.135 * len(text)
    shp = rect(slide, l, t, w, 0.4, WHITE, line=col, round_=True)
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = col; r.font.name = FONT


def new_slide(title, role, route, definition, features, point):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW.inches, SH.inches, WHITE)
    rect(s, 0, 0, SW.inches, 1.32, NAVY)
    rect(s, 0, 1.32, SW.inches, 0.05, MINT)
    textbox(s, 0.62, 0.26, 9.6, 0.7, [[(title, 26, WHITE, True, False)]])
    textbox(s, 0.64, 0.92, 9.6, 0.34, [[(route, 11.5, LTBLUE, False, False)]])
    col = ADMIN if role == "관리자" else NAVY
    role_chip_outline(s, 11.55, 0.42, role + " 화면", col)
    textbox(s, 12.2, 6.95, 0.9, 0.4, [[("00 / 00", 10, GRAY, False, False)]], align=PP_ALIGN.RIGHT)

    ph = rect(s, 0.4, 1.55, 7.5, 4.5, LIGHTER, line=BORDER, round_=True)
    tf = ph.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "화면 스크린샷 영역"; _set(r, 13, FAINT, bold=True)

    RIGHT_X, RIGHT_EDGE = 8.1, 12.9
    rw = RIGHT_EDGE - RIGHT_X
    textbox(s, RIGHT_X, 1.6, 2, 0.34, [[("정의", 12, MINT, True, False)]])
    textbox(s, RIGHT_X, 1.96, rw, 0.9, [[(definition, 13, GRAY, False, False)]])
    textbox(s, RIGHT_X, 2.86, 4, 0.34, [[("핵심 기능", 12, BLUE, True, False)]])
    fy = 3.28
    for head, body in features:
        rect(s, RIGHT_X + 0.04, fy + 0.06, 0.12, 0.12, BLUE)
        textbox(s, RIGHT_X + 0.33, fy - 0.04, rw - 0.33, 0.72,
                [[(head + "  ", 13, INK, True, False), (body, 11, GRAY, False, False)]])
        fy += 0.72
    box = rect(s, 0.62, 6.05, 12.1, 0.82, LIGHT, line=RGBColor(0xBA, 0xD6, 0xF4), round_=True)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "POINT  "; _set(r, 12, MINT, bold=True)
    r = p.add_run(); r.text = point; _set(r, 13, NAVY, bold=False)
    return s


NEW = ("상담 품질 검수 (관리자)", "관리자", "/post-consultation · audit-result",
       "센터 전체 상담을 AI 1차 검수 + 관리자 휴먼 2차 검수로 관리하는 품질 검수 화면입니다.",
       [("AI 1차 전수 검수", "모든 상담의 오안내·안내 누락을 자동 감지하고 위험도를 요약"),
        ("오안내 / 안내 누락 판정", "약관·내규 근거 + 매뉴얼 수행 여부로 정밀 판정"),
        ("휴먼 2차 검수", "관리자가 판정·코멘트를 작성하고 후속 조치를 지정"),
        ("검수 대기 큐 관리", "감지 건 중심으로 검수 대상을 우선 관리 → 상담사 이관")],
       "표본 점검을 넘어 전수 검수로 불완전판매·오안내 리스크를 사전에 관리합니다.")


def main():
    # 중복 방지
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip().startswith("상담 품질 검수 (관리자)"):
                print("이미 존재 — 스킵"); return

    before = len(prs.slides._sldIdLst)
    new_slide(*NEW)

    # 실시간 상담 모니터링(13, 0-based 12) 바로 뒤 = 0-based 13 에 삽입 → 14번 슬라이드
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    new_id = ids[-1]
    sldIdLst.remove(new_id)
    sldIdLst.insert(13, new_id)

    # 페이지 번호 재부여
    slides = list(prs.slides)
    total = len(slides)
    pat = re.compile(r"^\s*\d{1,2}\s*/\s*\d{1,2}\s*$")
    for pos, slide in enumerate(slides, 1):
        for sh in slide.shapes:
            if sh.has_text_frame and pat.match(sh.text_frame.text.strip()):
                sh.text_frame.paragraphs[0].runs[0].text = f"{pos:02d} / {total:02d}"
                break

    prs.save(FILE)
    print("saved:", FILE, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
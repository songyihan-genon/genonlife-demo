# -*- coding: utf-8 -*-
"""개요(2번) 리디자인 + 역할 색 체계 통일 + 3번~ 태그 아웃라인화. v3 → v4."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "GenOnLIFE_AICC_데모_v3.pptx"
OUT = "GenOnLIFE_AICC_데모_v4.pptx"

AGENT = RGBColor(0x0F, 0x34, 0x68)   # 상담사 — 네이비
ADMIN = RGBColor(0x0E, 0x9B, 0x86)   # 관리자 — 틸
COMMON = RGBColor(0x5B, 0x6B, 0x80)  # 공통 — 슬레이트
INK = RGBColor(0x10, 0x23, 0x3F)
GRAY = RGBColor(0x5B, 0x6B, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE6, 0xED, 0xF5)
FONT = "Apple SD Gothic Neo"
DIV_GRAD = [(0, "0F3468"), (45, "2F8BFF"), (100, "15C2A2")]

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height
slides = list(prs.slides)


def set_gradient(shape, stops, ang):
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = grad.makeelement(qn("a:gsLst"), {})
    for pos, rgb in stops:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": rgb}); gs.append(clr); gsLst.append(gs)
    grad.append(gsLst)
    grad.append(grad.makeelement(qn("a:lin"), {"ang": str(ang), "scaled": "1"}))
    ln = spPr.find(qn("a:ln"))
    (ln.addprevious(grad) if ln is not None else spPr.append(grad))


def role_color(text):
    a, m = "상담사" in text, "관리자" in text
    if a and m:
        return COMMON
    if a:
        return AGENT
    if m:
        return ADMIN
    return COMMON


# ════════ 1) 슬라이드 3번~ : 역할 태그를 아웃라인(흰 배경 + 테두리·글자색)으로 ════════
for slide in slides[2:]:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not (t.endswith("화면") and (sh.top or 0) < Emu(Inches(1.35))):
            continue
        col = role_color(t)
        sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = col; sh.line.width = Pt(1)
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = col


# ════════ 2) 개요(2번) 리디자인 ════════
def rect(slide, l, t, w, h, fill, line=None, round_=False, oval=False):
    shape = MSO_SHAPE.OVAL if oval else (MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(slide, l, t, w, h, runs, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        for (txt, size, color, bold) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT
    return tb


def role_chip(slide, l, t, text, fill):
    w = 0.28 + 0.135 * len(text)
    shp = rect(slide, l, t, w, 0.4, fill, round_=True)
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    return w


def menu_card(slide, x, y, w, h, name, desc, accent):
    rect(slide, x, y, w, h, WHITE, line=BORDER, round_=True)
    rect(slide, x + 0.26, y + 0.34, 0.13, 0.13, accent, oval=True)  # 역할 색 점
    textbox(slide, x + 0.5, y + 0.24, w - 0.65, 0.45, [[(name, 13.5, INK, True)]])
    textbox(slide, x + 0.27, y + 0.74, w - 0.5, h - 0.85, [[(desc, 10.5, GRAY, False)]])


ov = slides[1]
for sh in list(ov.shapes):
    sh._element.getparent().remove(sh._element)

rect(ov, 0, 0, SW.inches, SH.inches, WHITE)
top = rect(ov, 0, 0, SW.inches, 0.22, AGENT); set_gradient(top, DIV_GRAD, 0)
textbox(ov, 0.62, 0.52, 10, 0.7, [[("플랫폼 기능 요약", 28, INK, True)]])
textbox(ov, 0.64, 1.2, 11.8, 0.5,
        [[("상담 여정 전반을 메뉴별 AI 기능으로 지원합니다 — 상담사의 응대부터 관리자의 품질 운영까지", 13.5, GRAY, False)]])

AL, AR, GAP = 0.62, 12.71, 0.22

# 상담사
cwa = role_chip(ov, AL, 1.86, "상담사", AGENT)
textbox(ov, AL + cwa + 0.22, 1.9, 10, 0.4,
        [[("통화 중 실시간 보조와 자동 후처리로 응대 속도·정확도를 높입니다", 11.5, GRAY, False)]], anchor=MSO_ANCHOR.MIDDLE)
agent_cards = [
    ("AI 상담 홈", "콜 인입 수신 · 고객정보·예상 문의 자동 로딩 · 개인 업무 현황"),
    ("실시간 고객 상담", "STT · 상담 가이드 · 지식검색(RAG) · 업무정보 연동 · 관리자 코칭"),
    ("상담 이력 · 후처리", "이력 탐색 · 접촉이력 등록 · SMS 안내 발송 · 검수 결과 열람"),
    ("고객 민원 탐지", "STT 기반 불편 탐지 · 민원 발생 위험 평가 · VOC 자동 등록"),
]
n = len(agent_cards); cw = (AR - AL - GAP * (n - 1)) / n
for k, (nm, ds) in enumerate(agent_cards):
    menu_card(ov, AL + k * (cw + GAP), 2.34, cw, 1.72, nm, ds, AGENT)

# 관리자
cwm = role_chip(ov, AL, 4.34, "관리자", ADMIN)
textbox(ov, AL + cwm + 0.22, 4.38, 10, 0.4,
        [[("실시간 모니터링·코칭과 검수로 상담 품질을 상시 관리합니다", 11.5, GRAY, False)]], anchor=MSO_ANCHOR.MIDDLE)
admin_cards = [
    ("AI 상담 홈 (운영)", "실시간 운영 KPI · 상담사 상태·검수 분포·민원·품질 대시보드"),
    ("실시간 상담 모니터링", "관할 상담사 실시간 청취 · 고객정보·RAG 확인 · 실시간 코칭"),
    ("상담 검수", "AI 1차 검수 + 관리자 휴먼 2차 검수 · 코멘트·후속 조치"),
]
m = len(admin_cards); cw2 = (AR - AL - GAP * (m - 1)) / m
for k, (nm, ds) in enumerate(admin_cards):
    menu_card(ov, AL + k * (cw2 + GAP), 4.82, cw2, 1.72, nm, ds, ADMIN)

prs.save(OUT)
print("saved:", OUT)
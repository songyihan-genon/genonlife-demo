# -*- coding: utf-8 -*-
"""데모 PPT 스타일링: 타이틀 그라데이션 / 구분선 그라데이션 / 역할태그 흰색 / 기대효과 모던 리디자인. 새 파일로 저장."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "GenON_AICC_데모_.pptx"
OUT = "GenON_AICC_데모_v2.pptx"

NAVY = RGBColor(0x0F, 0x34, 0x68)
MINT = RGBColor(0x15, 0xC2, 0xA2)
INK = RGBColor(0x10, 0x23, 0x3F)
GRAY = RGBColor(0x5B, 0x6B, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE6, 0xED, 0xF5)
FONT = "Apple SD Gothic Neo"

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height


def inch(emu):
    return Emu(emu or 0).inches


def set_gradient(shape, stops, ang):
    """stops: [(pos0-100, 'RRGGBB')], ang: 60000ths deg (0=L→R)."""
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = grad.makeelement(qn("a:gsLst"), {})
    for pos, rgb in stops:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": rgb})
        gs.append(clr)
        gsLst.append(gs)
    grad.append(gsLst)
    lin = grad.makeelement(qn("a:lin"), {"ang": str(ang), "scaled": "1"})
    grad.append(lin)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)


TITLE_GRAD = [(0, "15457F"), (50, "0F3468"), (100, "0B2547")]   # 로그인 br 그라데이션
DIV_GRAD = [(0, "0F3468"), (45, "2F8BFF"), (100, "15C2A2")]      # 네이비→블루→민트(사이트 톤)


def white_text(shape):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = WHITE


# ── 슬라이드 순회 스타일 적용 ──
slides = list(prs.slides)
for si, slide in enumerate(slides, 1):
    title_text = ""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            title_text = sh.text_frame.text.strip().split("\n")[0]
            break

    for sh in list(slide.shapes):
        try:
            W, H, T = inch(sh.width), inch(sh.height), inch(sh.top)
        except Exception:
            continue
        # 1) 타이틀 슬라이드 풀배경 사각형 → 로그인 그라데이션
        if si == 1 and W > 12.5 and H > 6.5:
            set_gradient(sh, TITLE_GRAD, 2700000)  # 45° 대각
        # 2) 얇은 가로 구분선(민트 바) → 네이비-블루-민트 그라데이션
        elif H < 0.12 and W > 10:
            set_gradient(sh, DIV_GRAD, 0)
        # 3) 역할 태그(상단 '…화면' 칩) → 글자 흰색
        elif sh.has_text_frame and T < 1.35 and sh.text_frame.text.strip().endswith("화면"):
            white_text(sh)


# ── 기대 효과(마지막) 슬라이드 모던 리디자인 ──
def textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT
    return tb


def rect(slide, l, t, w, h, fill, line=None, round_=False, oval=False):
    shape = MSO_SHAPE.OVAL if oval else (MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


closing = slides[-1]
# 기존 도형 제거
for sh in list(closing.shapes):
    sh._element.getparent().remove(sh._element)

rect(closing, 0, 0, SW, SH, WHITE)
top = rect(closing, 0, 0, SW, Inches(0.22), MINT)
set_gradient(top, DIV_GRAD, 0)
textbox(closing, Inches(0.62), Inches(0.62), Inches(9), Inches(0.7), [[("기대 효과", 30, INK, True)]])
textbox(closing, Inches(0.64), Inches(1.32), Inches(11), Inches(0.5),
        [[("AICC 포털 도입으로 기대되는 운영·품질 변화", 14, GRAY, False)]])

EFFECTS = [
    ("응대 정확도", "실시간 STT·RAG·업무정보 연동으로 약관 기준과 고객 현황을 동시에 확인"),
    ("후처리 자동화", "상담 요약·SMS 초안·접촉이력 자동 생성으로 후처리 시간 단축"),
    ("선제 검수", "AI 1차 + 관리자 휴먼 검수로 불완전판매·오안내 리스크를 사전 관리"),
    ("민원 선제 탐지", "콜 STT 기반 불편 탐지·민원 위험 평가로 VOC를 사전 선별·등록"),
    ("실시간 코칭", "관리자가 상담 현장을 청취하며 즉시 코칭, 품질을 상시 관리"),
]
n = len(EFFECTS)
AL, AR, GAP = 0.62, 12.71, 0.22
CW = (AR - AL - GAP * (n - 1)) / n
CT, CH = 2.45, 3.9
for k, (head, body) in enumerate(EFFECTS):
    cx = AL + k * (CW + GAP)
    card = rect(closing, Inches(cx), Inches(CT), Inches(CW), Inches(CH), WHITE, line=BORDER, round_=True)
    # 번호 배지(그라데이션 원)
    badge = rect(closing, Inches(cx + 0.26), Inches(CT + 0.34), Inches(0.52), Inches(0.52), MINT, oval=True)
    set_gradient(badge, [(0, "0F3468"), (100, "15C2A2")], 2700000)
    btf = badge.text_frame; btf.word_wrap = False
    btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = f"{k+1:02d}"; br.font.size = Pt(13); br.font.bold = True; br.font.color.rgb = WHITE; br.font.name = FONT
    textbox(closing, Inches(cx + 0.24), Inches(CT + 1.05), Inches(CW - 0.48), Inches(0.6), [[(head, 15, INK, True)]])
    textbox(closing, Inches(cx + 0.24), Inches(CT + 1.6), Inches(CW - 0.48), Inches(2.1), [[(body, 11.5, GRAY, False)]])

prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))

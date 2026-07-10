# -*- coding: utf-8 -*-
"""portal-features.md 의 편집 가능한 PPTX 버전 (원본 디자인 충실 재현).
흰 배경 테마 · 네이비 제목+파란 밑줄 · 텍스트 좌측 / 스크린샷 우측 · 표 개요 · 섹션 간지.
모든 텍스트는 실제 텍스트박스(편집 가능), 스크린샷은 docs/img 이미지 삽입."""
import os, struct
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

IMG = "img"
NAVY = RGBColor(0x0F, 0x34, 0x68); BLUE = RGBColor(0x2F, 0x8B, 0xFF)
H3B = RGBColor(0x1A, 0x4F, 0x8F); TAGBG = RGBColor(0x0B, 0x4F, 0x91)
MINT = RGBColor(0x15, 0xC2, 0xA2)
INK = RGBColor(0x10, 0x23, 0x3F); GRAY = RGBColor(0x5B, 0x6B, 0x80)
NOTE = RGBColor(0x6B, 0x7F, 0x99); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCD, 0xD9, 0xE8); ROWALT = RGBColor(0xF2, 0xF8, 0xFF)
FAINT = RGBColor(0xA8, 0xB6, 0xC8)
FONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width.inches, prs.slide_height.inches


def _set(r, s, c, b=False, i=False):
    r.font.size = Pt(s); r.font.color.rgb = c; r.font.bold = b; r.font.italic = i; r.font.name = FONT


def rect(sl, l, t, w, h, fill, line=None, round_=False):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                              Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(sl, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, ital) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, ital)
    return tb


def tag_pill(sl, l, t, text):
    w = 0.24 + 0.098 * len(text)
    shp = rect(sl, l, t, w, 0.32, TAGBG, round_=True)
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.08); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set(r, 12, WHITE, True)
    return w


def png_size(path):
    d = open(path, "rb").read(24); return struct.unpack(">II", d[16:24])


def add_screenshot(sl, fname, bl=8.0, bt=1.55, bw=4.95, bh=4.6):
    path = os.path.join(IMG, fname)
    if not os.path.exists(path):
        rect(sl, bl, bt, bw, bh, ROWALT, line=BORDER, round_=True); return
    w, h = png_size(path); asp = w / h
    nw, nh = bw, bw / asp
    if nh > bh: nh, nw = bh, bh * asp
    pic = sl.shapes.add_picture(path, Inches(bl + (bw - nw) / 2), Inches(bt + (bh - nh) / 2),
                                width=Inches(nw), height=Inches(nh))
    pic.line.color.rgb = BORDER; pic.line.width = Pt(1)


# ─── 표지 (흰 배경, 중앙정렬) ───
def title_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    textbox(s, 1.0, 2.35, 11.33, 1.0, [[("제논라이프 AI 상담 포털", 40, NAVY, True, False)]], align=PP_ALIGN.CENTER)
    textbox(s, 1.0, 3.45, 11.33, 0.6, [[("기능 소개 — 상담사 · 관리자", 24, NAVY, True, False)]], align=PP_ALIGN.CENTER)
    rect(s, 5.17, 4.28, 3.0, 0.05, MINT)
    textbox(s, 1.0, 4.5, 11.33, 0.5,
            [[("콜센터 상담 지원 · VoC 통합 관리 · 상담 품질 검수 · 대외민원 대응", 15, GRAY, False, False)]], align=PP_ALIGN.CENTER)
    textbox(s, 1.0, 5.15, 11.33, 0.4, [[("AICC Demo · Next.js 14", 12, NOTE, False, True)]], align=PP_ALIGN.CENTER)


# ─── 섹션 간지 (흰 배경, 중앙) ───
def divider_slide(text):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    textbox(s, 1.0, 3.15, 11.33, 1.1, [[(text, 34, NAVY, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 5.67, 4.35, 2.0, 0.05, MINT)


# ─── 개요 (표) ───
def overview_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    textbox(s, 0.62, 0.5, 12, 0.6, [[("포털 개요 — 역할별 핵심 메뉴", 26, NAVY, True, False)]])
    rect(s, 0.62, 1.16, 12.1, 0.045, BLUE)
    textbox(s, 0.62, 1.28, 12.1, 0.5,
            [[("로그인 시 계정 유형(상담사 gena.kim / 관리자 park.admin)을 선택하면 사이드바 상단 메뉴가 분기됩니다.", 13, GRAY, False, False)]])
    rows = [
        ("", "상담사 (Agent)", "관리자 (Admin)"),
        ("1", "AI 상담 홈 (실시간 콜 허브)", "AI 상담 홈 — 운영 대시보드"),
        ("2", "실시간 고객 상담", "실시간 상담 모니터링"),
        ("3", "상담 이력 조회", "상담 품질 검수"),
        ("4", "후속업무 — 접촉 이력 등록", "VoC 애널리틱스 (실시간·통계·리포트)"),
        ("5", "후속업무 — SMS 발송", "민원 탐지·이관"),
        ("6", "후속업무 — 상담 검수 결과", "대외민원 처리"),
    ]
    tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.62), Inches(1.95), Inches(12.1), Inches(4.1))
    table = tbl_shape.table
    table.columns[0].width = Inches(0.7); table.columns[1].width = Inches(5.7); table.columns[2].width = Inches(5.7)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Inches(0.12); cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else ROWALT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            _set(r, 13.5, WHITE if ri == 0 else INK, (ri == 0 or ci == 0))
    textbox(s, 0.62, 6.25, 12.1, 0.5,
            [[("상담사는 \"상담 → 후처리 → 검수\" 현장 업무, 관리자는 \"운영 관제 → 품질 검수 → VoC·민원 대응\"을 담당합니다.", 12.5, NOTE, False, False)]])


# ─── 화면 슬라이드 (텍스트 좌 / 스크린샷 우) ───
def content_slide(idx, total, title, tag, purpose, sections, img):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    textbox(s, 0.62, 0.42, 12, 0.6, [[(title, 24, NAVY, True, False)]])
    rect(s, 0.62, 1.04, 12.1, 0.045, BLUE)
    # 태그 + 목적
    tw = tag_pill(s, 0.62, 1.2, tag)
    textbox(s, 0.62 + tw + 0.15, 1.19, 6.9 - tw - 0.15, 0.6, [[(purpose, 12, GRAY, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    # 스크린샷 우측
    add_screenshot(s, img)
    # 좌측 섹션(H3 + 불릿)
    y = 1.95
    for head, bullets in sections:
        textbox(s, 0.62, y, 7.1, 0.32, [[(head, 15, H3B, True, False)]])
        y += 0.36
        lines = [[("·  ", 11.5, BLUE, True, False), (b, 12, GRAY, False, False)] for b in bullets]
        tb = textbox(s, 0.62, y, 7.1, 0.36 * len(bullets) + 0.2, lines)
        for p in tb.text_frame.paragraphs:
            p.line_spacing = 1.06; p.space_after = Pt(3)
        y += 0.34 * len(bullets) + 0.22
    textbox(s, 12.15, 7.0, 0.95, 0.36, [[(f"{idx:02d} / {total:02d}", 10, GRAY, False, False)]], align=PP_ALIGN.RIGHT)


def closing_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    textbox(s, 1.0, 3.05, 11.33, 1.0, [[("감사합니다", 36, NAVY, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 5.67, 4.3, 2.0, 0.05, MINT)
    textbox(s, 1.0, 4.5, 11.33, 0.5,
            [[("제논라이프 AI 상담 포털 · 상담사 / 관리자 통합 지원", 15, GRAY, False, False)]], align=PP_ALIGN.CENTER)


AGENT = [
    ("[상담사] ① AI 상담 홈", "/main · agent",
     "근무 시작 시점의 실시간 콜 수신 대기 + 업무 지원 허브. 좌측 작업 존 / 우측 실시간 콜 패널.",
     [("오늘의 상담 KPI (상단 바)",
       ["상담 목표율 14/60 · 평균 응대대기(ASA) 0:21 · 평균 통화 5:18 · SLA 준수율 93% · 금주 만족도 4.6 · 후처리 대기 건수."]),
      ("나만의 업무 어시스턴트 (AI 검색·채팅형)",
       ["약관·규정·스크립트 자연어 검색 → 요약 답변.", "근거 태그(실손 특약 약관 §7 등) 클릭 → 원문 팝업."]),
      ("우측 실시간 콜 패널",
       ["콜 인입 자동 감지 → 고객 정보 조회(김민준·우량)·AI 예상 문의 표시 → \"통화 받기\".", "오늘의 고객 상담 이력 — 클릭 시 후처리로 이동."])],
     "agent-home.png"),

    ("[상담사] ② 실시간 고객 상담", "/insight-chat · counseling",
     "통화 수신 후 상담사를 실시간 보조하는 AI 상담 워크스페이스. 데모 케이스 김민준으로 시연.",
     [("상담 가이드 (좌측)",
       ["대화 키워드 실시간 추출 · 스크립트 가이드(도입·본인확인·동의).", "FCR 체크리스트 — 통화 녹취 고지·개인정보 동의·본인 확인 자동 체크 · 추천 응대 멘트."]),
      ("상담지식 어시스턴트 (RAG · 우측)",
       ["예시 칩·직접 질문 → 약관·업무기준 근거와 함께 답변. 관련 지식 자동 제시(FAQ)."]),
      ("관리자 실시간 코칭",
       ["상단 고객 배지(실효이력 3건·보험금 접수 중) · 관리자 코칭 메시지 실시간 수신."])],
     "agent-live.png"),

    ("[상담사] ③ 상담 이력 조회", "/post-consultation",
     "상담 이력을 기간·상태별로 검색·조회. 상담사는 본인 상담(김제나) 기준.",
     [("필터 바",
       ["보기: 통화별 / 상담사별 / 고객별 · 기간: 오늘 / 7일 / 30일 / 지정.", "AI 검수 · 휴먼 검수 상태 열 드롭다운 필터."]),
      ("이력 테이블",
       ["상담 ID(CL-20260514-042) · 고객 · 상담시간 · 채널(콜센터 IB·아웃바운드·챗봇 이관) · 주제 · 검수 상태.",
        "후속업무 열 — 접촉이력(등록완료/미등록) · SMS(발송완료/발송요청/미요청). 미처리는 노란 닷, 클릭 시 처리 화면 이동."]),
      ("데모 연동",
       ["오늘 데모: 접촉 미등록 5건 · SMS 미발송 3건이 후속업무와 연동."])],
     "agent-history.png"),

    ("[상담사] ④ 후속업무지원 (메일함형)", "?task=contact · sms · audit-result",
     "상담 종료 후 필수 후처리. 좌: 미처리 목록 / 우: 작업 패널.",
     [("접촉 이력 등록 (task=contact)",
       ["오늘 미등록 건 → 상담 개요 · 접촉 유형 분류(대·중·소) · 접촉 이력 3패널. AI가 유형별 초안 생성 → 검토·등록."]),
      ("SMS 발송 (task=sms)",
       ["미발송 건 → 문자 초안 어시스턴트(상담 분석→조회→근거·작성→요약 4단계) → 근거 기반 초안 검토·발송."]),
      ("상담 검수 결과 (task=audit-result)",
       ["AI 검수 결과 목록(감지 우선) → 상담 대화 + 검수 결과 2패널. 오안내·누락 확인 후 재안내/정정 조치."])],
     "agent-followup.png"),
]

ADMIN = [
    ("[관리자] ① AI 상담 홈 — 운영 대시보드", "/main · admin",
     "센터 전체를 한 화면에. 좌측 운영 대시보드 / 우측 4대 업무 큐.",
     [("콜상담 운영 관리",
       ["상담사 상태 도넛(상담중 16·후처리·대기·이석) · AI 1차 검수 도넛(통과 32·경미 8·심각 4).",
        "품질 지표(통과율 91%·오안내 3.2%·누락 5.1%·만족도 4.6) · 시간대별 콜 인입·응대(피크 14시 64콜)."]),
      ("VoC 처리 관리",
       ["채널별 인입(콜 182·이메일 74·챗봇 42·대외 14) · 부서 이관 상태 도넛(완료 248/76%)."]),
      ("우측 · 4대 업무 큐",
       ["실시간 상담 모니터링 · 상담 품질 검수 · 대외 민원 처리 · 민원 탐지 이관 — 긴급건 상단, 클릭 시 상세 이동."])],
     "admin-home.png"),

    ("[관리자] ② 실시간 상담 모니터링", "/realtime-monitoring",
     "진행 중인 통화를 실시간 청취·코칭하며 품질을 관제.",
     [("모니터링 로비",
       ["상담사 20명 카드(상담중 16·후처리 2·대기 1·이석 1), 상태·고객·주제·진행시간·AI 위험 신호.",
        "센터 필터 · 위험/이름 정렬 · 우측 KPI(평균 통화 4:18·오늘 처리 312건·만족도 94%)."]),
      ("모니터링 룸 (4분할)",
       ["실시간 STT 청취 + AI 신호 배너.", "실시간 코칭 — 예상 코칭 멘트 원클릭 전송(귓속말/전체).",
        "상담지식 RAG 검색 · 업무정보(고객·제품·청구) 패널."])],
     "admin-monitoring.png"),

    ("[관리자] ③ 상담 품질 검수", "?task=audit-result",
     "AI 1차 검수 + 관리자 휴먼 2차 검수 통합. 좌 목록 / 우 상세(메일함형).",
     [("검수 결과 목록",
       ["상태 점(심각·경미·통과) + 고객·상담ID·문제 요약. 감지 9건 / 통과 32건."]),
      ("검수 상세 (2열)",
       ["좌: 상담 대화 원문(STT 보정 이력 포함).", "우: 검수 결과 — AI 자동 판정(본인확인·안내 정확성·오안내·업무 누락) + 근거 약관 인용."]),
      ("휴먼 2차 검수",
       ["관리자 판정·코멘트 + 후속 조치(재안내 전화 / 정정 문자 / 접촉이력) → 상담사에게 이관, 완료 시 조치 이력 박제."])],
     "admin-audit.png"),

    ("[관리자] ④ VoC 애널리틱스 — Ⓐ 실시간 이슈 모니터링", "/voc-console · monitor",
     "실시간 VoC 스트림에서 위험 이슈를 조기 탐지.",
     [("실시간 이슈 알림 (카루셀)",
       ["\"보험금 부지급 불만 급증 +180%\", \"'금감원' 키워드 급증\", \"자동이체 이중 출금 다발\" 등."]),
      ("급증 키워드·유형",
       ["급증 키워드 TOP 10(5개씩 순환) · 키워드 클라우드.", "급증 유형 테이블 — 유형·주요 토픽·건수·전일 대비(보험금 부지급 1,300건 +24% 등)."]),
      ("부서별 처리 현황",
       ["6개 부서 처리/유입 · 처리율 바 · 상태(병목·주의·정상). 상단 채널·기간 필터."])],
     "voc-monitor.png"),

    ("[관리자] ④ VoC 애널리틱스 — Ⓑ 통계 대시보드", "statsboard",
     "전사 VoC를 매일 전일 데이터로 집계. 우측 AI 종합 분석 카드가 섹션별 인사이트 제공.",
     [("5개 섹션",
       ["① 일일 현황 — KPI 6종 스파크라인(전체 문의 1,850·VoC 750·고위험 95) / 구성 도넛 / 리스크 단계 퍼널.",
        "② 처리 실적 — 과녁 게이지(처리율 73%·SLA 88%·품질 85%) / 부서별 처리 현황(오늘 vs 전일).",
        "③ 기간 추세 — 탐지 추이(일별+시간대) / 인입 강도 히트맵(요일×시간대).",
        "④ 유형 분석 — 고객 프로파일 레이더 / 유형 포지셔닝 버블 / 상품별 점유율 트리맵 / 원인 분류.",
        "⑤ 인입 플로우 — 채널→유형→부서 생키(기본 접힘)."])],
     "voc-stats.png"),

    ("[관리자] ④ VoC 애널리틱스 — Ⓒ VoC 리포트", "report",
     "일일·월간 VoC 리포트를 자동 생성·시각화·배포.",
     [("리포트 생성 설정",
       ["일일 / 월간 토글 · 보고 일자·기간 · 센터·채널 선택 · 포함 섹션 다중 선택. 우측 생성 이력 패널."]),
      ("리포트 본문",
       ["접수 추이 라인 · 고객 경험 분포 도넛 · 금감원 민원 미니 바.",
        "주요 VoC 접수 내용 테이블 · 경험별 현황 · (월간) VoC 현황 표·상위 유형·개선 사례."]),
      ("내보내기",
       ["PDF · 엑셀 · 인쇄 · 공유."])],
     "voc-report.png"),

    ("[관리자] ⑤ 민원 탐지·이관 — VoC 통합 분석", "/complaint-detection",
     "콜·이메일·챗봇 문의를 단일 인입 큐로 통합(탐지·평가 → 유형 분류 → 부서 이관). 금일 인입 2,450건.",
     [("현황 대시보드 (통합/실시간/부서별 탭)",
       ["VoC 비중 ↔ 악성민원 발전 위험 · 채널별/유형별 고위험 비중 · 이관 진행률 73% · 부서별 처리 현황."]),
      ("인박스 테이블 (미배정 / 배정완료 / 처리완료)",
       ["위험도 점 · 고객·유형 · 접수일시 · 담당 부서 · 긴급도 · 상태. 부서 재배정·유형 변경 가능."]),
      ("상세 패널",
       ["처리 스테퍼(이관 대기→배정→처리) · 감지 키워드·트리거 · AI 분류·부서 배정 · VoC 등록."])],
     "voc-detect.png"),

    ("[관리자] ⑥ 대외민원 처리", "/external-complaint",
     "금융감독원 등 대외기관 이첩 민원 회신 초안 작성·처리. 3단 레이아웃.",
     [("좌 · 금감원 접수함",
       ["미처리/완료/전체 · 처리 구분(이첩·자율조정·사실조회) · D-Day 임박 배지."]),
      ("중앙 · 민원 내용 + 초안 작성",
       ["원천 시스템 조회(계약원장·청구심사·손해사정·CRM) · 사실관계 타임라인.",
        "금감원 회신 ↔ 민원인 회신 탭 · 템플릿 선택 · 자동 초안 생성 · 검증 체크리스트(근거·금지표현·기한)."]),
      ("우 · 인용 자료",
       ["사실 / 근거(약관·법령) / 사례(유사 분쟁조정 판례) → 초안에 자동 첨부."])],
     "external.png"),
]


def build():
    title_slide()
    overview_slide()
    total = len(AGENT) + len(ADMIN)
    idx = 1
    divider_slide("Ⅰ. 상담사 (Agent) 화면")
    for sl in AGENT:
        content_slide(idx, total, *sl); idx += 1
    divider_slide("Ⅱ. 관리자 (Admin) 화면")
    for sl in ADMIN:
        content_slide(idx, total, *sl); idx += 1
    closing_slide()
    out = "portal-features-editable.pptx"
    prs.save(out)
    print("saved:", out, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
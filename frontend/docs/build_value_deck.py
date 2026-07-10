# -*- coding: utf-8 -*-
"""GenON LIFE AICC 포털 — 가치 중심 화면 소개 덱.
'정의/핵심기능 나열' 대신 [가치 헤드라인 + 가치 3블록(가치→실현 방식) + POINT] 포맷.
스크린샷 영역은 빈 플레이스홀더. build_deck/relayout/style2 색·밴드 계승."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x0F, 0x34, 0x68)
BLUE = RGBColor(0x2F, 0x6B, 0xB0)
MINT = RGBColor(0x15, 0xC2, 0xA2)
ADMIN = RGBColor(0x0E, 0x9B, 0x86)
COMMON = RGBColor(0x5B, 0x6B, 0x80)
LIGHT = RGBColor(0xF2, 0xF8, 0xFF)
LIGHTER = RGBColor(0xF7, 0xFA, 0xFE)
INK = RGBColor(0x10, 0x23, 0x3F)
GRAY = RGBColor(0x5B, 0x6B, 0x80)
FAINT = RGBColor(0xA8, 0xB6, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCD, 0xD9, 0xE8)
CARDBORDER = RGBColor(0xE6, 0xED, 0xF5)
LTBLUE = RGBColor(0x9F, 0xC4, 0xEE)
FONT = "Apple SD Gothic Neo"
DIV_GRAD = [(0, "0F3468"), (45, "2F8BFF"), (100, "15C2A2")]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width.inches, prs.slide_height.inches


def role_color(role):
    return {"상담사": NAVY, "관리자": ADMIN, "공통": COMMON}.get(role, NAVY)


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = FONT


def set_gradient(shape, stops, ang):
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = grad.makeelement(qn("a:gsLst"), {})
    for pos, rgb in stops:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": rgb}); gs.append(clr); gsLst.append(gs)
    grad.append(gsLst)
    grad.append(grad.makeelement(qn("a:lin"), {"ang": str(ang), "scaled": "1"}))
    ln = spPr.find(qn("a:ln"))
    (ln.addprevious(grad) if ln is not None else spPr.append(grad))


def rect(s, l, t, w, h, fill, line=None, round_=False, oval=False):
    shape = MSO_SHAPE.OVAL if oval else (MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE)
    shp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, italic) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, italic)
    return tb


def chip(s, l, t, text, fill, fg, outline=False):
    w = 0.28 + 0.135 * len(text)
    shp = rect(s, l, t, w, 0.4, WHITE if outline else fill, line=fill if outline else None, round_=True)
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set(r, 13, fill if outline else fg, bold=True)
    return w


# ═══════════════════ 표지 ═══════════════════
def title_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 3.45, SW, 0.06, MINT)
    textbox(s, 0.9, 1.95, 11.5, 1.2,
            [[("GenON ", 46, WHITE, True, False), ("LIFE", 46, MINT, True, False),
              ("  ·  AICC 포털", 46, WHITE, True, False)]])
    textbox(s, 0.92, 3.62, 11.5, 0.9,
            [[("AI 상담 어시스턴트가 만드는 변화 — 화면별 핵심 가치", 20, RGBColor(0xCF, 0xE0, 0xF1), False, False)]])
    textbox(s, 0.92, 4.42, 11.8, 0.9,
            [[("응대 정확도 · 후처리 자동화 · 선제 검수 · 민원 선제 탐지 · 실시간 코칭", 14, LTBLUE, False, False)]])
    textbox(s, 0.9, 6.7, 11.5, 0.5,
            [[("© 2026 GenON LIFE · AICC Portal 데모", 11, RGBColor(0x6F, 0x8C, 0xB5), False, False)]])


# ═══════════════════ 개요(플랫폼 기능 요약) ═══════════════════
def menu_card(s, x, y, w, h, name, desc, accent):
    rect(s, x, y, w, h, WHITE, line=CARDBORDER, round_=True)
    rect(s, x + 0.24, y + 0.28, 0.12, 0.12, accent, oval=True)
    textbox(s, x + 0.46, y + 0.18, w - 0.6, 0.45, [[(name, 12.5, INK, True, False)]])
    textbox(s, x + 0.26, y + 0.62, w - 0.48, h - 0.72, [[(desc, 9.8, GRAY, False, False)]])


def overview_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    top = rect(s, 0, 0, SW, 0.22, NAVY); set_gradient(top, DIV_GRAD, 0)
    textbox(s, 0.62, 0.46, 10, 0.7, [[("플랫폼 기능 요약", 27, INK, True, False)]])
    textbox(s, 0.64, 1.12, 11.9, 0.5,
            [[("상담 여정 전반을 메뉴별 AI 가치로 지원합니다 — 상담사의 응대부터 관리자의 품질·민원 운영까지", 13, GRAY, False, False)]])
    AL, AR, GAP = 0.62, 12.71, 0.2

    chip(s, AL, 1.74, "상담사", NAVY, WHITE)
    agent = [
        ("AI 상담 홈", "콜 인입 자동 로딩 · 개인 KPI · 업무 어시스턴트 검색"),
        ("실시간 고객 상담", "STT · 상담 가이드 · 지식검색(RAG) · 관리자 코칭"),
        ("상담 이력 · 후처리", "이력 탐색 · 접촉이력 · SMS 안내 · 검수 피드백"),
        ("상담 검수 피드백", "AI 검수·관리자 코멘트 수신 → 재발 방지"),
    ]
    n = len(agent); cw = (AR - AL - GAP * (n - 1)) / n
    for k, (nm, ds) in enumerate(agent):
        menu_card(s, AL + k * (cw + GAP), 2.18, cw, 1.42, nm, ds, BLUE)

    chip(s, AL, 3.86, "관리자", ADMIN, WHITE)
    admin = [
        ("AI 상담 홈 (운영)", "운영 KPI · 상담사 상태·검수 분포·품질 대시보드"),
        ("실시간 상담 모니터링", "관할 상담사 실시간 청취 · 실시간 코칭"),
        ("상담 검수", "AI 1차 + 관리자 휴먼 2차 검수 · 후속 조치"),
        ("VoC 애널리틱스", "실시간 이슈 · 통계 대시보드 · VoC 리포트"),
        ("민원 탐지·이관", "채널 통합 큐 · AI 분류 · 부서 이관"),
        ("대외민원 처리", "금감원 회신 초안 · 근거·사례 자동 인용"),
    ]
    m = 3; cw2 = (AR - AL - GAP * (m - 1)) / m
    for k, (nm, ds) in enumerate(admin):
        row, col = divmod(k, m)
        menu_card(s, AL + col * (cw2 + GAP), 4.3 + row * 1.5, cw2, 1.36, nm, ds, ADMIN)


# ═══════════════════ 가치 중심 화면 슬라이드 ═══════════════════
def value_slide(idx, total, title, role, route, headline, values, point, note=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, SW, 1.32, NAVY)
    rect(s, 0, 1.32, SW, 0.05, MINT)
    textbox(s, 0.62, 0.26, 9.8, 0.7, [[(title, 25, WHITE, True, False)]])
    textbox(s, 0.64, 0.94, 9.8, 0.34, [[(route, 11.5, LTBLUE, False, False)]])
    rc = role_color(role)
    chip(s, 11.5, 0.42, role + " 화면", rc, WHITE, outline=True)
    textbox(s, 12.2, 6.98, 0.9, 0.36, [[(f"{idx:02d} / {total:02d}", 10, GRAY, False, False)]], align=PP_ALIGN.RIGHT)

    # 좌측 스크린샷 플레이스홀더
    ph = rect(s, 0.4, 1.62, 6.5, 4.42, LIGHTER, line=BORDER, round_=True)
    tf = ph.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "화면 스크린샷 영역"; _set(r, 13, FAINT, bold=True)
    if note:
        textbox(s, 0.4, 6.02, 6.5, 0.3, [[(note, 10, FAINT, False, True)]], align=PP_ALIGN.CENTER)

    # 우측 상단 — 가치 헤드라인
    RX = 7.2
    textbox(s, RX, 1.6, 5.6, 0.3, [[("이 화면이 만드는 변화", 11.5, MINT, True, False)]])
    textbox(s, RX, 1.95, 5.65, 1.15, [[(headline, 18, NAVY, True, False)]])

    # 가치 3블록 (번호 배지 + 가치명 + 실현 방식)
    y = 3.28
    for i, (vname, how) in enumerate(values, 1):
        badge = rect(s, RX, y, 0.4, 0.4, rc, oval=True)
        bt = badge.text_frame; bt.word_wrap = False; bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bt.margin_top = bt.margin_bottom = bt.margin_left = bt.margin_right = 0
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = str(i); _set(br, 14, WHITE, bold=True)
        textbox(s, RX + 0.56, y - 0.06, 5.1, 0.4, [[(vname, 14, INK, True, False)]])
        textbox(s, RX + 0.56, y + 0.34, 5.15, 0.62, [[(how, 11, GRAY, False, False)]])
        y += 0.92

    # POINT
    box = rect(s, 0.62, 6.05, 12.1, 0.82, LIGHT, line=RGBColor(0xBA, 0xD6, 0xF4), round_=True)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "POINT  "; _set(r, 12, MINT, bold=True)
    r = p.add_run(); r.text = point; _set(r, 13, NAVY, bold=False)


# ═══════════════════ 클로징(기대 효과) ═══════════════════
def closing_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 2.05, SW, 0.05, MINT)
    textbox(s, 0.9, 1.15, 11.5, 0.8, [[("기대 효과", 30, WHITE, True, False)]])
    textbox(s, 0.92, 1.62, 11.5, 0.4,
            [[("AICC 포털 도입으로 기대되는 운영·품질 변화", 13.5, LTBLUE, False, False)]])
    items = [
        ("응대 정확도", "실시간 STT·RAG·업무정보 연동으로 약관 기준과 고객 실제 현황을 동시에 확인"),
        ("후처리 자동화", "상담 요약·SMS 초안·접촉이력을 자동 생성해 후처리 시간 단축"),
        ("선제 검수", "AI 1차 + 관리자 휴먼 검수로 불완전판매·오안내 리스크를 사전 관리"),
        ("민원 선제 탐지", "콜·멀티채널 VoC를 탐지·정량화해 민원 확산 전에 선별·이관"),
        ("실시간 코칭", "관리자가 상담 현장을 청취하며 즉시 코칭, 품질을 상시 관리"),
    ]
    y = 2.45
    for i, (head, body) in enumerate(items, 1):
        badge = rect(s, 0.95, y, 0.42, 0.42, MINT, oval=True)
        bt = badge.text_frame; bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bt.margin_top = bt.margin_bottom = bt.margin_left = bt.margin_right = 0
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = f"0{i}"; _set(br, 13, NAVY, bold=True)
        textbox(s, 1.55, y - 0.04, 11.2, 0.7,
                [[(head + "   ", 16, WHITE, True, False), (body, 12.5, RGBColor(0xCF, 0xE0, 0xF1), False, False)]])
        y += 0.86


# ═══════════════════ 콘텐츠 정의 ═══════════════════
SLIDES = [
    ("로그인 · 역할 선택", "공통", "/login",
     "하나의 포털, 역할에 맞춘 두 개의 업무 경험",
     [("역할 최적화", "상담사/관리자 선택 시 메뉴·홈·권한이 자동으로 재구성"),
      ("빠른 진입", "계정 유형 선택 + 보안 로그인으로 즉시 업무 시작"),
      ("일관된 브랜드 경험", "AICC 포털 통합 진입 화면으로 통일된 UX 제공")],
     "상담 현장과 운영 관리를 하나의 포털에서 역할에 맞춰 제공합니다."),

    ("AI 상담 홈 (상담사)", "상담사", "/main",
     "콜을 받기 전에, AI가 이미 응대 준비를 끝냅니다",
     [("응대 준비 자동화", "콜 인입 시 고객 정보·예상 문의를 자동으로 로딩"),
      ("즉시 지식 접근", "약관·규정·스크립트를 자연어로 검색하고 근거 원문까지 확인"),
      ("업무 장악력", "개인 KPI·후처리 대기 건을 한 화면에서 파악·진입")],
     "상담 시작 전 준비 시간을 없애 응대 속도와 정확도를 동시에 끌어올립니다."),

    ("실시간 고객 상담 · 콜 대기", "상담사", "/insight-chat · counseling",
     "비어 있던 대기 시간을 상담 품질을 높이는 학습 시간으로",
     [("상시 업스킬", "대기 중 약관·업무 기준·상품 카탈로그를 학습"),
      ("응대 표준화", "유형별 FCR 필수 안내 체크리스트를 미리 숙지"),
      ("끊김 없는 전환", "콜 연결 시 실시간 STT 상담 화면으로 자동 전환")],
     "유휴 시간을 상담 역량을 끌어올리는 준비 시간으로 전환합니다."),

    ("실시간 고객 상담 · 상담 시작", "상담사", "/insight-chat · counseling",
     "약관을 찾는 상담이 아니라, 고객에게 집중하는 상담",
     [("응대 정확도 향상", "실시간 STT 분석 → 약관·기준(RAG)과 고객 현황을 동시에 제시"),
      ("응대 속도 향상", "대화 유형별 FCR 필수 안내·추천 스크립트를 자동 제시"),
      ("리스크 최소화", "민감 구간에서 관리자가 접속해 실시간으로 코칭·개입")],
     "RAG는 '약관 기준'을, 업무정보는 '고객 실제 현황'을 분담해 빈틈없이 안내합니다.",
     "※ 본 데모는 '김민준' 고객 케이스에 한해 동작합니다."),

    ("상담 이력 조회", "상담사·관리자", "/post-consultation",
     "상담 이후의 모든 후속 업무를 놓치지 않는 한 흐름",
     [("빠른 추적", "기간·검수 조건으로 이력 탐색 → 요약·원문·코칭·검수 상세 조회"),
      ("누락 방지", "미등록·발송요청·검수 감지 건을 시각적으로 표시"),
      ("즉시 처리 연결", "해당 셀 클릭 시 접촉이력·SMS·검수 화면으로 바로 이동")],
     "처리할 후속 업무를 이력에서 곧장 연결·완료해 누락을 없앱니다."),

    ("접촉이력 등록", "상담사", "/post-consultation · contact",
     "상담사가 직접 쓰던 접촉이력을, AI가 초안까지 작성",
     [("입력 부담 절감", "종료된 상담을 분석해 유형(대/중/소)별 요약 초안을 자동 생성"),
      ("정확한 분류", "발화지점별 다중 의도 분석으로 유형 판단 근거를 제시"),
      ("맥락 이해", "고객의 과거 접촉 이력을 타임라인으로 함께 제공")],
     "AI가 접촉이력을 초안화하고 상담사는 검토·확정만 — 후처리 부담을 줄입니다.",
     "※ 본 데모는 '김민준' 고객 케이스에 한해 동작합니다."),

    ("SMS 안내 작성", "상담사", "/post-consultation · sms",
     "상담 요약을 그대로 고객 안내 문자로 연결",
     [("초안 자동 생성", "상담 문의 내용을 자동 정리·근거 탐색 후 안내 문구 초안 작성"),
      ("자연어 편집", "원하는 톤·문구를 대화형으로 요청하면 즉시 반영"),
      ("오안내 방지", "필수 고지·금지 표현을 자동 점검해 안내 누락 차단")],
     "근거 기반 안내 문자를 자동 생성해 후속 안내를 빠르고 정확하게 마무리합니다.",
     "※ 본 데모는 '김민준' 고객 케이스에 한해 동작합니다."),

    ("상담 검수 결과 (상담사)", "상담사", "/post-consultation · audit-result",
     "검수 결과가 상담사의 다음 응대를 바꿉니다",
     [("명확한 피드백", "AI 오안내·안내 누락 감지 결과와 관리자 코멘트를 함께 수신"),
      ("근거 기반 이해", "판정 사유 + 약관·내규 근거 + 매뉴얼 수행 여부를 제시"),
      ("즉각 개선", "고객 재안내 등 후속 조치를 그 자리에서 수행")],
     "AI 감지 → 관리자 코멘트 → 후속 조치로 이어지는 검수 루프로 품질을 지속 개선합니다.",
     "※ 본 데모는 '정해린' 고객 케이스에 한해 동작합니다."),

    ("AI 상담 홈 (관리자)", "관리자", "/main",
     "센터 전체의 운영과 품질을 한 화면에서 장악",
     [("실시간 현황 파악", "상담 중·대기 콜·오안내/누락 감지 현황을 KPI로 확인"),
      ("품질 가시성", "상담사 상태·검수 분포·민원 통계·품질 지표를 시각화"),
      ("즉시 조치", "실시간 모니터링·검수 대기 큐로 필요한 작업에 바로 진입")],
     "운영·품질·민원을 하나의 대시보드로 통합해 관리자의 의사결정을 가속합니다."),

    ("실시간 상담 모니터링", "관리자", "/realtime-monitoring",
     "녹취를 나중에 듣는 관리가 아니라, 지금 개입하는 관리",
     [("실시간 관제", "관할 상담사 선택 → 위험·주의 신호와 실시간 STT를 청취"),
      ("즉시 코칭", "귓속말/전체 안내로 상담 현장에서 바로 코칭을 전달"),
      ("동일 시야 확보", "고객 업무정보·RAG를 상담사와 같은 화면으로 확인")],
     "사후 피드백의 지연을 없애고 상담 현장에서 실시간으로 품질을 관리합니다."),

    ("상담 검수 (관리자)", "관리자", "/post-consultation · audit-result",
     "표본 점검을 넘어, 전수 검수로 리스크를 선제 관리",
     [("전수 검수", "AI 1차 검수로 모든 상담의 오안내·안내 누락을 감지"),
      ("정밀 판정", "오안내/누락 탭 + 약관·내규 근거로 관리자가 휴먼 2차 검수"),
      ("검수 루프 종결", "코멘트·후속 조치를 상담사에게 이관해 재발을 방지")],
     "불완전판매·오안내 리스크를 사후가 아닌 사전에 관리합니다."),

    ("VoC 애널리틱스 · 실시간 이슈", "관리자", "/voc-console · monitor",
     "민원이 터지기 전에, 이상 징후를 먼저 포착",
     [("조기 경보", "부지급 불만·'금감원' 키워드 급증 등 실시간 이슈 알림"),
      ("위험 가시화", "급증 키워드 TOP·유형별 건수·전일 대비 추이 표시"),
      ("병목 관제", "부서별 처리/유입·처리율·병목 상태를 실시간 모니터링")],
     "이상 징후를 실시간으로 감지해 민원이 확산되기 전에 선제 대응합니다."),

    ("VoC 애널리틱스 · 통계 대시보드", "관리자", "/voc-console · statsboard",
     "숫자만 보여주는 대시보드가 아니라, 해석까지 주는 대시보드",
     [("전사 집계", "문의·VoC·고위험 KPI, 구성·리스크 단계, 처리 실적을 한눈에"),
      ("원인 분석", "유형·상품·고객 프로파일·인입 플로우로 근본 원인 파악"),
      ("AI 인사이트", "섹션마다 위험 요인·병목·권고 사항을 자동 요약")],
     "흩어진 VoC 지표를 집계하고 AI가 해석·권고까지 덧붙입니다."),

    ("VoC 애널리틱스 · VoC 리포트", "관리자", "/voc-console · report",
     "반복 작성하던 VoC 리포트를, 클릭 몇 번으로",
     [("자동 편성", "접수 추이·고객 경험·주요 VoC·개선 사례를 문서형으로 구성"),
      ("규제 대응", "금감원 민원 제기·처리 집계를 리포트에 포함"),
      ("즉시 배포", "일일/월간 생성 후 PDF·엑셀·인쇄·공유")],
     "정기 VoC 리포트를 표준 서식으로 자동 생성해 보고 부담을 줄입니다."),

    ("민원 탐지·이관", "관리자", "/complaint-detection",
     "흩어진 채널의 문의를, 하나의 큐에서 분류·이관",
     [("채널 통합 인입", "콜·이메일·챗봇 문의를 단일 인입 큐로 통합"),
      ("AI 분류", "유형 자동 분류·부서 추천 → 재배정·유형 변경"),
      ("신속 이관", "위험도·긴급도 기반 우선순위로 담당 부서 배정·등록")],
     "채널별로 흩어진 문의를 한 큐로 모아 분류·평가·부서 이관을 일괄 처리합니다."),

    ("대외민원 처리", "관리자", "/external-complaint",
     "대외기관 회신을, 사실 조회부터 근거 인용까지 표준화",
     [("사실관계 자동 집계", "계약원장·청구심사·손해사정·CRM 원천 시스템을 자동 조회"),
      ("근거 기반 초안", "회신 템플릿 자동 초안 + 약관·법령·유사 판례 인용"),
      ("리스크 차단", "근거 인용·금지 표현·회신 기한을 검증 체크리스트로 관리")],
     "사실 조회 → 근거 인용 → 초안 작성을 표준 절차로 지원해 회신 품질을 높입니다."),
]


def build():
    title_slide()
    overview_slide()
    total = len(SLIDES) + 4  # 표지·개요·클로징 제외한 페이지 번호 기준(대략)
    for i, sl in enumerate(SLIDES, start=3):
        value_slide(i, len(SLIDES) + 2, *sl)
    closing_slide()
    out = "GenOnLIFE_AICC_데모_가치중심.pptx"
    prs.save(out)
    print("saved:", out, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
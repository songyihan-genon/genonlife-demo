# -*- coding: utf-8 -*-
"""GenON LIFE AICC 포털 데모 — 화면별 핵심기능 PPT 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0F, 0x34, 0x68)
BLUE = RGBColor(0x2F, 0x6B, 0xB0)
MINT = RGBColor(0x15, 0xC2, 0xA2)
LIGHT = RGBColor(0xF2, 0xF8, 0xFF)
LIGHTER = RGBColor(0xF7, 0xFA, 0xFE)
INK = RGBColor(0x10, 0x23, 0x3F)
GRAY = RGBColor(0x5B, 0x6B, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
FONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT


def textbox(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold, italic) or list of such for one paragraph with multiple runs."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        runs = para if isinstance(para, list) else [para]
        for (txt, size, color, bold, italic) in runs:
            r = p.add_run()
            r.text = txt
            _set(r, size, color, bold, italic)
    return tb


def rect(slide, l, t, w, h, fill, line=None, round_=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def chip(slide, l, t, text, fill, fg):
    w = Inches(0.18 + 0.105 * len(text))
    h = Inches(0.34)
    shp = rect(slide, l, t, w, h, fill, round_=True)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set(r, 11, fg, bold=True)
    return w


def title_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    # accent bar
    rect(s, 0, Inches(3.45), SW, Inches(0.06), MINT)
    textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4), [
        [("GenON ", 46, WHITE, True, False), ("LIFE", 46, MINT, True, False), ("  ·  AICC 포털", 46, WHITE, True, False)],
    ])
    textbox(s, Inches(0.92), Inches(3.65), Inches(11.5), Inches(0.9), [
        [("AI 기반 컨택센터 상담 어시스턴트 — 화면별 핵심 기능", 20, RGBColor(0xCF, 0xE0, 0xF1), False, False)],
    ])
    textbox(s, Inches(0.92), Inches(4.45), Inches(11.5), Inches(0.9), [
        [("실시간 상담 보조 · 후처리 자동화 · 오안내·누락 검수 · 민원 탐지 · 실시간 코칭", 14, RGBColor(0x9F, 0xC4, 0xEE), False, False)],
    ])
    textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5), [
        [("© 2026 GenON LIFE · AICC Portal 데모", 11, RGBColor(0x6F, 0x8C, 0xB5), False, False)],
    ])


def section_label(s, text):
    textbox(s, Inches(0.62), Inches(1.62), Inches(4), Inches(0.4),
            [[("정의", 11, MINT, True, False)]])


def content_slide(idx, total, title, role, route, definition, features, highlight):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    # top band
    rect(s, 0, 0, SW, Inches(1.32), NAVY)
    rect(s, 0, Inches(1.32), SW, Inches(0.05), MINT)
    textbox(s, Inches(0.62), Inches(0.26), Inches(9.6), Inches(0.7),
            [[(title, 26, WHITE, True, False)]])
    textbox(s, Inches(0.64), Inches(0.92), Inches(9.6), Inches(0.34),
            [[(route, 11.5, RGBColor(0x9F, 0xC4, 0xEE), False, False)]])
    # role chip (top-right)
    rc_fill = MINT if role == "관리자" else RGBColor(0x3D, 0xB0, 0xFF)
    chip(s, Inches(11.55), Inches(0.42), role + " 화면", rc_fill, NAVY)
    # page number
    textbox(s, Inches(12.2), Inches(6.95), Inches(0.9), Inches(0.4),
            [[(f"{idx:02d} / {total:02d}", 10, GRAY, False, False)]], align=PP_ALIGN.RIGHT)

    # 정의
    textbox(s, Inches(0.62), Inches(1.62), Inches(2), Inches(0.34), [[("정의", 12, MINT, True, False)]])
    textbox(s, Inches(0.62), Inches(1.96), Inches(12.1), Inches(0.7),
            [[(definition, 14.5, GRAY, False, False)]])

    # 핵심 기능
    textbox(s, Inches(0.62), Inches(2.78), Inches(4), Inches(0.34), [[("핵심 기능", 12, BLUE, True, False)]])
    fy = Inches(3.2)
    for head, body in features:
        # bullet square
        rect(s, Inches(0.66), fy + Emu(int(Inches(0.06))), Inches(0.12), Inches(0.12), BLUE)
        textbox(s, Inches(0.95), fy - Inches(0.04), Inches(11.6), Inches(0.62),
                [[(head + "  ", 14.5, INK, True, False), (body, 13, GRAY, False, False)]])
        fy = fy + Inches(0.62)

    # highlight box
    hb_t = Inches(6.05)
    box = rect(s, Inches(0.62), hb_t, Inches(12.1), Inches(0.82), LIGHT, line=RGBColor(0xBA, 0xD6, 0xF4), round_=True)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "POINT  "; _set(r, 12, MINT, bold=True)
    r = p.add_run(); r.text = highlight; _set(r, 13, NAVY, bold=False)


def closing_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, Inches(2.1), SW, Inches(0.05), MINT)
    textbox(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8), [[("기대 효과", 30, WHITE, True, False)]])
    items = [
        ("응대 정확도 향상", "실시간 STT·RAG·업무정보 연동으로 약관 기준과 고객 실제 현황을 동시에 확인"),
        ("후처리 자동화", "상담 요약·SMS 초안·접촉이력을 자동 생성해 후처리 시간 단축"),
        ("오안내·누락 선제 검수", "AI 1차 검수 + 관리자 휴먼 검수로 불완전판매·오안내 리스크 관리"),
        ("민원 선제 탐지", "콜 STT 기반 불편 탐지·민원 위험 평가로 VOC를 사전 선별·등록"),
        ("실시간 코칭 체계", "관리자가 상담 현장을 청취하며 즉시 코칭, 품질을 상시 관리"),
    ]
    y = Inches(2.5)
    for head, body in items:
        rect(s, Inches(0.95), y + Inches(0.07), Inches(0.14), Inches(0.14), MINT)
        textbox(s, Inches(1.25), y - Inches(0.04), Inches(11.2), Inches(0.7),
                [[(head + "   ", 16, WHITE, True, False), (body, 13, RGBColor(0xCF, 0xE0, 0xF1), False, False)]])
        y = y + Inches(0.78)


SLIDES = [
    ("로그인 · 역할 선택", "공통", "/login",
     "상담사·관리자 계정 유형을 선택해 진입하는 역할 기반 로그인 화면입니다.",
     [("계정 유형 선택", "상담사 / 관리자 2개 역할로 진입"),
      ("역할별 화면 분기", "좌측 메뉴·홈 레이아웃·권한이 역할에 따라 자동 구성"),
      ("브랜드 진입 경험", "AICC 포털 브랜드 패널 + 보안 로그인 폼")],
     "하나의 포털에서 상담사 업무와 관리자 운영 화면을 역할에 맞춰 제공합니다."),

    ("AI 상담 홈 (상담사)", "상담사", "/main",
     "상담사의 업무 시작점 — 실시간 콜 수신과 개인 업무 대시보드를 제공합니다.",
     [("콜 인입 패널", "실시간 콜 수신 → 통화 받기, 고객 정보·예상 문의 자동 로딩"),
      ("나만의 업무 어시스턴트", "약관·규정·스크립트를 자연어로 빠르게 검색"),
      ("바로가기 · 오늘의 상담 이력", "자주 쓰는 업무 화면 + 후처리 필요 건 바로 진입"),
      ("개인 KPI", "오늘의 상담 목표율·후처리 대기 등 업무 현황")],
     "콜을 받기 전에 고객 이력·예상 문의까지 AI가 미리 준비해 응대 속도를 높입니다."),

    ("실시간 고객 상담 (상담사)", "상담사", "/insight-chat · counseling",
     "통화 중 실시간으로 상담사를 보조하는 AICC 핵심 화면입니다.",
     [("실시간 STT", "통화를 실시간 받아쓰고 발화에서 키워드를 자동 추출"),
      ("상담 가이드", "FCR 필수 안내 자동 체크 · 대화 유형별 추천 스크립트"),
      ("상담지식 어시스턴트 (RAG)", "약관·업무 기준 지식 질의응답 + 출처 제시"),
      ("업무정보 연동", "계약·처리계 데이터(제출 현황·처리 가능·예상 환급금)"),
      ("관리자 실시간 코칭", "민감 구간에서 관리자가 접속해 실시간 코칭")],
     "RAG는 '약관·기준'을, 업무정보는 '고객의 실제 현황'을 분담해 정확히 안내합니다."),

    ("실시간 고객 상담 · 콜 대기 화면", "상담사", "/insight-chat · counseling",
     "통화가 없는 대기 시간에 업무를 숙지하는 기본(대기) 화면입니다.",
     [("콜 대기 상태", "콜 연결 시 실시간 STT가 시작됨을 안내"),
      ("필수·자주 상담 가이드", "도입·본인확인부터 해지·민원 응대까지 단계별 가이드"),
      ("제논라이프 상품 카탈로그", "전체 판매 상품을 유형별로 확인"),
      ("상담지식 검색", "빈 시간에 약관·업무 기준을 미리 학습")],
     "유휴 시간을 상담 품질을 끌어올리는 학습 시간으로 전환합니다."),

    ("SMS 안내 작성", "상담사", "/post-consultation · sms",
     "상담 내용을 근거로 고객 안내 문자를 생성·발송하는 화면입니다.",
     [("근거 기반 초안 생성", "어시스턴트 근거 → 안내 문구 초안 자동 작성"),
      ("상담사 편집", "초안을 검토·수정 후 발송"),
      ("발송 처리", "발송 완료 표기 + 표준 안내 문구")],
     "상담 요약을 그대로 고객 안내 문자로 연결해 후속 안내를 자동화합니다."),

    ("상담 이력 조회", "상담사·관리자", "/post-consultation",
     "콜 이력을 탐색하고 후속업무·검수 상태를 한눈에 보는 화면입니다.",
     [("기간·검수 필터", "최근 7일/30일, 검수 대기/완료 등 조건 탐색"),
      ("후속업무 상태 연동", "접촉이력·SMS 처리 상태를 실시간 반영"),
      ("셀 바로가기", "미등록·발송요청·검수 감지 셀 클릭 시 해당 화면으로 즉시 이동"),
      ("상세 카드", "상담 요약·대화 원문·관리자 코칭 이력·검수 결과")],
     "상담 이후의 모든 후속 업무를 이력 한 화면에서 추적·처리합니다."),

    ("접촉이력 등록", "상담사", "/post-consultation · contact",
     "종료된 상담을 분석해 접촉 유형을 분류하고 이력으로 등록합니다.",
     [("자동 분석", "종료 상담을 분석해 요약·유형 분류 초안 생성"),
      ("유형별 요약 편집", "접촉 유형(대/중/소)별 요약 초안을 검토·수정"),
      ("Tele-Pro 이력 등록", "표준 형식으로 접촉이력 등록 → 상담 개요에 반영")],
     "상담사가 직접 작성하던 접촉이력을 AI가 초안화해 입력 부담을 줄입니다."),

    ("상담 검수 결과", "관리자", "/post-consultation · audit-result",
     "AI 1차 검수와 관리자 휴먼 2차 검수를 한 화면에 담은 검수 리포트입니다.",
     [("AI 검수 종합", "오안내·안내 누락 감지 결과와 위험도 요약"),
      ("오안내 / 안내 누락 탭", "판정·평가 사유 + 약관·내규 근거 + 매뉴얼 수행 여부"),
      ("휴먼 검수 (관리자)", "검수 판정 선택 · 코멘트 직접 작성 · 후속 조치 선택"),
      ("상담사 피드백", "검토 결과를 상담사에게 전달해 재발 방지")],
     "AI 감지 → 관리자 판정·코멘트 → 후속 조치까지 검수 루프를 한 화면에서 닫습니다."),

    ("고객 민원 탐지", "상담사·관리자", "/complaint-detection",
     "콜 STT 기반으로 불편을 탐지하고 민원 위험을 평가해 VOC로 등록합니다.",
     [("인사이트 대시보드", "민원 위험 추이 · 상품 유형별 민원 · 상담 주의 포인트"),
      ("잠재 민원 케이스", "불편 신호·민원 발전 위험·제도 개선 영역을 케이스별로 정리"),
      ("AI 평가 기준 관리", "불편/위험/긴급/분류 기준을 수정·추가, 적용 가능성 검토"),
      ("VOC 자동 등록", "위험도 기반 선별 → 표준 형식 등록(기존 접수 제외·실패 재수행)")],
     "흩어진 불편 신호를 AI가 선별·정량화해 민원으로 번지기 전에 대응합니다."),

    ("AI 상담 홈 (관리자)", "관리자", "/main",
     "관할 상담사 운영 현황과 상담 품질을 보는 관리자 대시보드입니다.",
     [("운영 KPI", "실시간 상담 중·대기 콜·오안내/누락 감지 현황"),
      ("시각화 대시보드", "상담사 상태·검수 결과 분포·민원 통계·품질 지표·추이"),
      ("실시간 상담 모니터링", "상담 중 상담사 입장(청취) / 메시지 발송"),
      ("검수 대기 큐", "상담번호·담당 상담사 중심으로 검수 대상 관리")],
     "상담사 화면과 같은 자리에서 관리자에게 필요한 운영·품질 관점을 제공합니다."),

    ("실시간 상담 모니터링 (관리자)", "관리자", "/realtime-monitoring",
     "실시간 코칭을 위한 관할 상담사 상담 모니터링 화면입니다.",
     [("관할 상담사 선택", "상담 중 상담사 선택 → 위험·주의 신호 표시"),
      ("실시간 STT 청취", "선택 상담사의 실시간 상담 내용을 읽기 전용으로 모니터링"),
      ("고객 정보·RAG", "고객 업무정보와 약관·업무 기준을 함께 확인"),
      ("실시간 코칭 채팅", "상담사에게 즉시 코칭 전달(귓속말/전체 안내)")],
     "관리자가 상담 현장을 실시간으로 보며 그 자리에서 코칭하는 메인 작업 화면입니다."),
]


def build():
    title_slide()
    total = len(SLIDES) + 2  # title + closing
    for i, sl in enumerate(SLIDES, start=2):
        content_slide(i, total, *sl)
    closing_slide()
    out = "GenON_AICC_데모.pptx"
    prs.save(out)
    print("saved:", out, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
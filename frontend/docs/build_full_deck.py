# -*- coding: utf-8 -*-
"""GenON LIFE AICC 포털 — 화면별 핵심 기능 덱(전체 클린 빌드).
모든 슬라이드를 python-pptx로 생성(빈 스크린샷 플레이스홀더) → 모든 렌더러 호환.
포맷: 정의 / 핵심 기능(라벨+설명) / POINT + 좌측 스크린샷 영역. 0616 덮어쓰기."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x0F, 0x34, 0x68); BLUE = RGBColor(0x2F, 0x6B, 0xB0)
MINT = RGBColor(0x15, 0xC2, 0xA2); ADMIN = RGBColor(0x0E, 0x9B, 0x86)
COMMON = RGBColor(0x5B, 0x6B, 0x80)
LIGHT = RGBColor(0xF2, 0xF8, 0xFF); LIGHTER = RGBColor(0xF7, 0xFA, 0xFE)
INK = RGBColor(0x10, 0x23, 0x3F); GRAY = RGBColor(0x5B, 0x6B, 0x80)
FAINT = RGBColor(0xA8, 0xB6, 0xC8); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCD, 0xD9, 0xE8); CARDBORDER = RGBColor(0xE6, 0xED, 0xF5)
LTBLUE = RGBColor(0x9F, 0xC4, 0xEE); PTBORDER = RGBColor(0xBA, 0xD6, 0xF4)
FONT = "Apple SD Gothic Neo"
DIV_GRAD = [(0, "0F3468"), (45, "2F8BFF"), (100, "15C2A2")]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width.inches, prs.slide_height.inches


def rc(role):
    return {"상담사": NAVY, "관리자": ADMIN, "공통": COMMON, "상담사·관리자": COMMON}.get(role, NAVY)


def _set(r, s, c, b=False, i=False):
    r.font.size = Pt(s); r.font.color.rgb = c; r.font.bold = b; r.font.italic = i; r.font.name = FONT


def set_gradient(shape, stops, ang):
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {}); gsLst = grad.makeelement(qn("a:gsLst"), {})
    for pos, rgb in stops:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": rgb}); gs.append(clr); gsLst.append(gs)
    grad.append(gsLst); grad.append(grad.makeelement(qn("a:lin"), {"ang": str(ang), "scaled": "1"}))
    ln = spPr.find(qn("a:ln")); (ln.addprevious(grad) if ln is not None else spPr.append(grad))


def rect(sl, l, t, w, h, fill, line=None, round_=False, oval=False):
    shape = MSO_SHAPE.OVAL if oval else (MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE)
    shp = sl.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(sl, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, ital) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, ital)
    return tb


def chip(sl, l, t, text, col, outline=False):
    w = 0.28 + 0.135 * len(text)
    shp = rect(sl, l, t, w, 0.4, WHITE if outline else col, line=col if outline else None, round_=True)
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set(r, 13, col if outline else WHITE, True)


# ─────── 표지 ───────
def title_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 3.45, SW, 0.06, MINT)
    textbox(s, 0.9, 1.95, 11.5, 1.2,
            [[("GenON ", 46, WHITE, True, False), ("LIFE", 46, MINT, True, False),
              ("  ·  AICC 포털", 46, WHITE, True, False)]])
    textbox(s, 0.92, 3.62, 11.5, 0.9,
            [[("AI 기반 컨택센터 상담 어시스턴트 — 화면별 핵심 기능", 20, RGBColor(0xCF, 0xE0, 0xF1), False, False)]])
    textbox(s, 0.92, 4.42, 11.8, 0.9,
            [[("실시간 상담 보조 · 후처리 자동화 · 오안내·누락 검수 · 민원 탐지 · 실시간 코칭", 14, LTBLUE, False, False)]])
    textbox(s, 0.9, 6.7, 11.5, 0.5,
            [[("© 2026 GenON LIFE · AICC Portal 데모", 11, RGBColor(0x6F, 0x8C, 0xB5), False, False)]])


# ─────── 개요 ───────
def menu_card(s, x, y, w, h, name, desc, accent):
    rect(s, x, y, w, h, WHITE, line=CARDBORDER, round_=True)
    rect(s, x + 0.24, y + 0.28, 0.12, 0.12, accent, oval=True)
    textbox(s, x + 0.46, y + 0.18, w - 0.6, 0.45, [[(name, 12.5, INK, True, False)]])
    textbox(s, x + 0.26, y + 0.62, w - 0.48, h - 0.72, [[(desc, 9.8, GRAY, False, False)]])


def overview_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    top = rect(s, 0, 0, SW, 0.22, NAVY); set_gradient(top, DIV_GRAD, 0)
    textbox(s, 0.62, 0.46, 10, 0.7, [[("플랫폼 기능 요약", 27, INK, True, False)]])
    textbox(s, 0.64, 1.12, 11.9, 0.5,
            [[("상담 여정 전반을 메뉴별 AI 기능으로 지원합니다 — 상담사의 응대부터 관리자의 품질·민원 운영까지", 13, GRAY, False, False)]])
    AL, AR, GAP = 0.62, 12.71, 0.2
    chip(s, AL, 1.74, "상담사", NAVY)
    agent = [
        ("AI 상담 홈", "콜 인입 자동 로딩 · 개인 KPI · 업무 어시스턴트 검색"),
        ("실시간 고객 상담", "STT · 상담 가이드 · 지식검색(RAG) · 관리자 코칭"),
        ("상담 이력 · 후처리", "이력 탐색 · 접촉이력 · SMS 안내 · 검수 피드백"),
        ("상담 검수 피드백", "AI 검수·관리자 코멘트 수신 → 재발 방지"),
    ]
    n = len(agent); cw = (AR - AL - GAP * (n - 1)) / n
    for k, (nm, ds) in enumerate(agent):
        menu_card(s, AL + k * (cw + GAP), 2.18, cw, 1.42, nm, ds, BLUE)
    chip(s, AL, 3.86, "관리자", ADMIN)
    admin = [
        ("AI 상담 홈 (운영)", "운영 KPI · 상담사 상태·검수 분포·품질 대시보드"),
        ("실시간 상담 모니터링", "관할 상담사 실시간 청취 · 실시간 코칭"),
        ("상담 품질 검수", "AI 1차 + 관리자 휴먼 2차 검수 · 후속 조치"),
        ("VoC 애널리틱스", "실시간 이슈 · 통계 대시보드 · VoC 리포트"),
        ("민원 탐지·이관", "채널 통합 큐 · AI 분류 · 부서 이관"),
        ("대외민원 처리", "금감원 회신 초안 · 근거·사례 자동 인용"),
    ]
    m = 3; cw2 = (AR - AL - GAP * (m - 1)) / m
    for k, (nm, ds) in enumerate(admin):
        row, col = divmod(k, m)
        menu_card(s, AL + col * (cw2 + GAP), 4.3 + row * 1.5, cw2, 1.36, nm, ds, ADMIN)


# ─────── 화면 슬라이드(정의/핵심기능/POINT) ───────
def content_slide(idx, total, title, role, route, definition, features, point, note=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, SW, 1.32, NAVY)
    rect(s, 0, 1.32, SW, 0.05, MINT)
    textbox(s, 0.62, 0.26, 9.6, 0.7, [[(title, 26, WHITE, True, False)]])
    textbox(s, 0.64, 0.94, 9.6, 0.34, [[(route, 11.5, LTBLUE, False, False)]])
    chip(s, 11.5, 0.42, role + " 화면", rc(role), outline=True)
    textbox(s, 12.15, 6.98, 0.95, 0.36, [[(f"{idx:02d} / {total:02d}", 10, GRAY, False, False)]], align=PP_ALIGN.RIGHT)
    # 좌측 스크린샷 영역
    ph = rect(s, 0.4, 1.55, 7.5, 4.5, LIGHTER, line=BORDER, round_=True)
    tf = ph.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "화면 스크린샷 영역"; _set(r, 13, FAINT, True)
    if note:
        textbox(s, 0.4, 6.02, 7.5, 0.3, [[(note, 9.5, FAINT, False, True)]], align=PP_ALIGN.CENTER)
    # 우측 정의/핵심기능/POINT
    RX, RE = 8.1, 12.9; rw = RE - RX
    textbox(s, RX, 1.6, 2, 0.34, [[("정의", 12, MINT, True, False)]])
    textbox(s, RX, 1.96, rw, 0.9, [[(definition, 13, GRAY, False, False)]])
    textbox(s, RX, 2.86, 4, 0.34, [[("핵심 기능", 12, BLUE, True, False)]])
    step = 0.72 if len(features) >= 4 else 0.86
    fy = 3.28 if len(features) >= 4 else 3.4
    for head, body in features:
        rect(s, RX + 0.04, fy + 0.06, 0.12, 0.12, BLUE)
        textbox(s, RX + 0.33, fy - 0.04, rw - 0.33, step,
                [[(head + "  ", 13, INK, True, False), (body, 11, GRAY, False, False)]])
        fy += step
    box = rect(s, 0.62, 6.05, 12.1, 0.82, LIGHT, line=PTBORDER, round_=True)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "POINT  "; _set(r, 12, MINT, True)
    r = p.add_run(); r.text = point; _set(r, 13, NAVY, False)


# ─────── 클로징(기대 효과) ───────
def closing_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 2.05, SW, 0.05, MINT)
    textbox(s, 0.9, 1.15, 11.5, 0.8, [[("기대 효과", 30, WHITE, True, False)]])
    textbox(s, 0.92, 1.62, 11.5, 0.4, [[("AICC 포털 도입으로 기대되는 운영·품질 변화", 13.5, LTBLUE, False, False)]])
    items = [
        ("응대 정확도", "실시간 STT·RAG·업무정보 연동으로 약관 기준과 고객 실제 현황을 동시에 확인"),
        ("후처리 자동화", "상담 요약·SMS 초안·접촉이력을 자동 생성해 후처리 시간 단축"),
        ("선제 검수", "AI 1차 + 관리자 휴먼 검수로 불완전판매·오안내 리스크를 사전 관리"),
        ("민원 선제 탐지", "콜·멀티채널 VoC를 탐지·정량화해 민원 확산 전에 선별·이관"),
        ("실시간 코칭", "관리자가 상담 현장을 청취하며 즉시 코칭, 품질을 상시 관리"),
    ]
    y = 2.45
    for i, (head, body) in enumerate(items, 1):
        badge = rect(s, 0.95, y, 0.42, 0.42, MINT, oval=True)
        bt = badge.text_frame; bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bt.margin_top = bt.margin_bottom = bt.margin_left = bt.margin_right = 0
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = f"0{i}"; _set(br, 13, NAVY, True)
        textbox(s, 1.55, y - 0.04, 11.2, 0.7,
                [[(head + "   ", 16, WHITE, True, False), (body, 12.5, RGBColor(0xCF, 0xE0, 0xF1), False, False)]])
        y += 0.86


SLIDES = [
    ("로그인 · 역할 선택", "공통", "/login",
     "상담사·관리자 계정 유형을 선택해 진입하는 역할 기반 로그인 화면입니다.",
     [("계정 유형 선택", "상담사 / 관리자 2개 역할로 진입"),
      ("역할별 화면 분기", "좌측 메뉴·홈 레이아웃·권한이 역할에 따라 자동 구성"),
      ("브랜드 진입 경험", "AICC 포털 브랜드 패널 + 보안 로그인 폼")],
     "하나의 포털에서 상담사 업무와 관리자 운영 화면을 역할에 맞춰 제공합니다."),

    ("AI 상담 홈 (상담사)", "상담사", "/main",
     "상담사의 업무 시작점 — 실시간 콜 수신과 개인 업무 대시보드를 제공합니다.",
     [("콜 인입 패널", "실시간 콜 수신 → 통화 받기, 고객 정보·예상 문의 자동 로딩"),
      ("나만의 업무 어시스턴트", "약관·규정·스크립트를 자연어로 빠르게 검색"),
      ("바로가기 · 오늘의 상담 이력", "자주 쓰는 업무 화면 + 후처리 필요 건 바로 진입"),
      ("개인 KPI", "오늘의 상담 목표율·후처리 대기 등 업무 현황")],
     "콜을 받기 전에 고객 이력·예상 문의까지 AI가 미리 준비해 응대 속도를 높입니다."),

    ("실시간 고객 상담 · 콜 대기", "상담사", "/insight-chat · counseling",
     "통화가 없는 대기 시간에 업무를 숙지하는 기본(대기) 화면입니다.",
     [("콜 대기 상태", "콜 연결 시 실시간 STT가 시작됨을 안내"),
      ("필수·자주 상담 가이드", "도입·본인확인부터 해지·민원 응대까지 단계별 가이드"),
      ("제논라이프 상품 카탈로그", "전체 판매 상품을 유형별로 확인"),
      ("상담지식 검색", "빈 시간에 약관·업무 기준을 미리 학습")],
     "유휴 시간을 상담 품질을 끌어올리는 학습 시간으로 전환합니다."),

    ("실시간 고객 상담 · 상담 시작", "상담사", "/insight-chat · counseling",
     "통화 중 실시간으로 상담사를 보조하는 AICC 핵심 화면입니다.",
     [("실시간 STT", "통화를 실시간 받아쓰고 발화에서 키워드를 자동 추출"),
      ("상담 가이드", "FCR 필수 안내 자동 체크 · 대화 유형별 추천 스크립트"),
      ("상담지식 어시스턴트 (RAG)", "약관·업무 기준 질의응답 + 출처 제시"),
      ("관리자 실시간 코칭", "민감 구간에서 관리자가 접속해 실시간 코칭")],
     "RAG는 '약관·기준'을, 업무정보는 '고객의 실제 현황'을 분담해 정확히 안내합니다.",
     "※ 본 데모는 '김민준' 고객 케이스에 한해 동작합니다."),

    ("SMS 안내 작성", "상담사", "/post-consultation · sms",
     "상담 내용을 근거로 고객 안내 문자를 생성·발송하는 화면입니다.",
     [("근거 기반 초안 생성", "상담 문의 내용을 자동 정리·근거 탐색 후 안내 문구 초안 작성"),
      ("대화형 수정", "원하는 톤·문구를 자연어로 요청하면 즉시 반영"),
      ("필수 고지 자동 점검", "필수 고지·금지 표현을 자동 검수해 오안내·안내 누락 방지")],
     "상담 요약을 그대로 고객 안내 문자로 연결해 후속 안내를 자동화합니다.",
     "※ 본 데모는 '김민준' 고객 케이스에 한해 동작합니다."),

    ("접촉이력 등록", "상담사", "/post-consultation · contact",
     "종료된 상담을 분석해 접촉 유형을 분류하고 이력으로 등록합니다.",
     [("과거 접촉이력", "고객 과거 접촉 이력을 타임라인으로 표시해 맥락 이해"),
      ("발화지점별 다중의도 분석", "발화 구간별 의도를 분석해 검토 용이"),
      ("유형별 요약 편집", "접촉 유형(대/중/소)별 요약 초안 검토·수정·등록")],
     "상담사가 직접 작성하던 접촉이력을 AI가 초안화해 입력 부담을 줄입니다.",
     "※ 본 데모는 '김민준' 고객 케이스에 한해 동작합니다."),

    ("상담 이력 조회", "상담사·관리자", "/post-consultation",
     "콜 이력을 탐색하고 후속업무·검수 상태를 한눈에 보는 화면입니다.",
     [("기간·검수 필터", "최근 7일/30일, 검수 대기/완료 등 조건 탐색"),
      ("후속업무 상태 연동", "접촉이력·SMS 처리 상태를 실시간 반영"),
      ("셀 바로가기", "미등록·발송요청·검수 감지 셀 클릭 시 해당 화면으로 즉시 이동"),
      ("상세 카드", "상담 요약·대화 원문·관리자 코칭 이력·검수 결과")],
     "상담 이후의 모든 후속 업무를 이력 한 화면에서 추적·처리합니다."),

    ("상담 검수 결과 (상담사)", "상담사", "/post-consultation · audit-result",
     "AI 검수 결과와 관리자 코멘트를 피드백으로 확인하고 후속 조치하는 상담사 화면입니다.",
     [("AI 검수 종합", "오안내·안내 누락 감지 결과와 위험도 요약"),
      ("오안내 / 안내 누락 탭", "판정·평가 사유 + 약관·내규 근거 + 매뉴얼 수행 여부"),
      ("관리자 코멘트·피드백", "관리자 2차 검수 코멘트를 피드백으로 수신"),
      ("후속 조치", "고객 재안내 등 피드백 기반 후속 조치 수행")],
     "AI 감지 → 관리자 코멘트 → 후속 조치로 이어지는 검수 루프로 응대 품질을 개선합니다.",
     "※ 본 데모는 '정해린' 고객 케이스에 한해 동작합니다."),

    ("AI 상담 홈 (관리자)", "관리자", "/main",
     "관할 상담사 운영 현황과 상담 품질을 보는 관리자 대시보드입니다.",
     [("운영 KPI", "실시간 상담 중·대기 콜·오안내/누락 감지 현황"),
      ("시각화 대시보드", "상담사 상태·검수 결과 분포·민원 통계·품질 지표·추이"),
      ("실시간 상담 모니터링", "상담 중 상담사 입장(청취) / 메시지 발송"),
      ("검수 대기 큐", "AI 1차 검수 오류 감지 건 중심으로 검수 대상 관리")],
     "운영·품질·민원을 한 대시보드로 통합해 관리자의 의사결정을 가속합니다."),

    ("실시간 상담 모니터링 (관리자)", "관리자", "/realtime-monitoring",
     "실시간 코칭을 위한 관할 상담사 상담 모니터링 화면입니다.",
     [("관할 상담사 선택", "상담 중 상담사 선택 → AI 위험·주의 신호 표시"),
      ("실시간 STT 청취", "선택 상담사의 실시간 상담 내용을 읽기 전용으로 모니터링"),
      ("고객 정보·RAG", "고객 업무정보와 약관·업무 기준을 함께 확인"),
      ("실시간 코칭 채팅", "상담사에게 즉시 코칭 전달(귓속말/전체 안내)")],
     "관리자가 상담 현장을 실시간으로 보며 그 자리에서 코칭하는 메인 작업 화면입니다."),

    ("상담 품질 검수 (관리자)", "관리자", "/post-consultation · audit-result",
     "센터 전체 상담을 AI 1차 검수 + 관리자 휴먼 2차 검수로 관리하는 품질 검수 화면입니다.",
     [("AI 1차 전수 검수", "모든 상담의 오안내·안내 누락을 자동 감지하고 위험도를 요약"),
      ("오안내 / 안내 누락 판정", "약관·내규 근거 + 매뉴얼 수행 여부로 정밀 판정"),
      ("휴먼 2차 검수", "관리자가 판정·코멘트를 작성하고 후속 조치를 지정"),
      ("검수 대기 큐 관리", "감지 건 중심으로 검수 대상을 우선 관리 → 상담사 이관")],
     "표본 점검을 넘어 전수 검수로 불완전판매·오안내 리스크를 사전에 관리합니다."),

    ("VoC 애널리틱스 · 실시간 이슈 모니터링", "관리자", "/voc-console · monitor",
     "실시간 VoC 인입에서 급증 이슈·위험 키워드를 조기에 탐지하는 관제 화면입니다.",
     [("실시간 이슈 알림", "부지급 불만·'금감원' 키워드 급증 등 이상 징후를 카루셀 경보로 표시"),
      ("급증 키워드·유형", "키워드 TOP 10·클라우드 + 급증 VoC 유형(건수·전일 대비)"),
      ("부서별 처리 현황", "6개 부서 처리/유입·처리율과 병목·주의·정상 상태"),
      ("채널·기간 필터", "콜센터·이메일·대외기관 채널 및 기간별로 재집계")],
     "이상 징후를 실시간으로 감지해 민원이 확산되기 전에 선제 대응합니다."),

    ("VoC 애널리틱스 · 통계 대시보드", "관리자", "/voc-console · statsboard",
     "전사 VoC를 매일 전일 데이터로 집계하고, AI가 섹션별 인사이트를 제시하는 통계 화면입니다.",
     [("일일 현황", "문의·VoC·고위험 KPI(스파크라인) · VoC 구성 도넛 · 리스크 단계 퍼널"),
      ("처리 실적", "처리율·SLA·품질 게이지 + 부서별 처리 현황(오늘 vs 전일)"),
      ("기간 추세·유형 분석", "탐지 추이·인입 강도 히트맵 · 고객 프로파일·상품별 점유율"),
      ("AI 종합 분석", "섹션마다 위험 요인·병목·권고 사항을 자동 요약")],
     "흩어진 VoC 지표를 하나의 대시보드로 집계하고 AI가 해석까지 덧붙입니다."),

    ("VoC 애널리틱스 · VoC 리포트", "관리자", "/voc-console · report",
     "일일·월간 VoC 리포트를 자동 생성·시각화하고 배포하는 화면입니다.",
     [("리포트 설정", "일일/월간 · 기간·센터·채널 · 포함 섹션 선택"),
      ("자동 구성", "접수 추이·고객 경험 분포·주요 VoC·개선 사례를 문서형으로 편성"),
      ("금감원 민원 집계", "제기·처리 건수와 주요 민원 요약 포함"),
      ("내보내기", "PDF · 엑셀 · 인쇄 · 공유")],
     "반복 작성하던 VoC 리포트를 클릭 몇 번으로 표준 서식에 맞춰 생성합니다."),

    ("민원 탐지·이관", "관리자", "/complaint-detection",
     "콜·이메일·챗봇 문의를 단일 인입 큐로 통합해 유형 분류부터 부서 이관까지 처리하는 화면입니다.",
     [("통합 현황 대시보드", "VoC 비중·고위험 비중·이관 진행률·부서별 처리 현황"),
      ("통합 인입 큐", "미배정/배정완료/처리완료 탭 · 위험도·긴급도·담당 부서"),
      ("AI 분류·부서 이관", "유형 자동 분류·부서 추천 → 재배정·유형 변경"),
      ("VoC 등록·처리", "감지 키워드·트리거 확인 후 표준 형식으로 등록")],
     "흩어진 채널의 문의를 한 큐로 모아 분류·평가·부서 이관을 일괄 처리합니다."),

    ("대외민원 처리", "관리자", "/external-complaint",
     "금융감독원 등 대외기관 이첩 민원의 회신 초안을 작성·검증·처리하는 화면입니다.",
     [("금감원 접수함", "미처리·기한 임박(D-Day) 우선 정렬 · 처리 구분(이첩·자율조정·사실조회)"),
      ("원천 시스템 조회", "계약원장·청구심사·손해사정·CRM 사실관계 자동 집계"),
      ("자동 초안 작성", "금감원·민원인 회신 템플릿 + 검증 체크리스트(근거·금지표현·기한)"),
      ("근거·사례 첨부", "약관·법령·유사 분쟁조정 판례를 초안에 자동 인용")],
     "사실 조회 → 근거 인용 → 초안 작성을 표준 절차로 지원해 회신 품질을 높입니다."),
]


def build():
    title_slide()
    overview_slide()
    total = len(SLIDES) + 3  # 표지·개요·클로징 포함 전체 페이지 수
    for i, sl in enumerate(SLIDES, start=3):
        content_slide(i, total, *sl)
    closing_slide()
    out = "GenOnLIFE_AICC_데모_0616.pptx"
    prs.save(out)
    print("saved:", out, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
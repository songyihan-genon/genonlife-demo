# -*- coding: utf-8 -*-
"""0615 덱에 새로 추가된 관리자 화면(VoC 애널리틱스 3종·민원 탐지·이관·대외민원 처리)을
build_deck / relayout / style2 와 동일한 스타일로 슬라이드 추가.
스크린샷 영역은 빈 플레이스홀더로 남겨둠. 기대 효과 슬라이드 앞에 삽입."""
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = "GenOnLIFE_AICC_데모_0615.pptx"
OUT = "GenOnLIFE_AICC_데모_0616.pptx"

NAVY = RGBColor(0x0F, 0x34, 0x68)
BLUE = RGBColor(0x2F, 0x6B, 0xB0)
MINT = RGBColor(0x15, 0xC2, 0xA2)
ADMIN = RGBColor(0x0E, 0x9B, 0x86)   # 관리자 — 틸 (style2 기준)
LIGHT = RGBColor(0xF2, 0xF8, 0xFF)
LIGHTER = RGBColor(0xF7, 0xFA, 0xFE)
INK = RGBColor(0x10, 0x23, 0x3F)
GRAY = RGBColor(0x5B, 0x6B, 0x80)
FAINT = RGBColor(0xA8, 0xB6, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCD, 0xD9, 0xE8)
LTBLUE = RGBColor(0x9F, 0xC4, 0xEE)
FONT = "Apple SD Gothic Neo"

# relayout 좌측 이미지 영역
IMG_L, IMG_T, IMG_W, IMG_H = 0.4, 1.55, 7.5, 4.5
# 우측 텍스트 영역
RIGHT_X, RIGHT_EDGE = 8.1, 12.9

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = FONT


def rect(slide, l, t, w, h, fill, line=None, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, italic) in (para if isinstance(para, list) else [para]):
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, italic)
    return tb


def role_chip_outline(slide, l, t, text, col):
    """style2 아웃라인 칩 — 흰 배경 + 색 테두리·글자."""
    w = 0.28 + 0.135 * len(text)
    shp = rect(slide, l, t, w, 0.4, WHITE, line=col, round_=True)
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = col; r.font.name = FONT


def new_slide(title, role, route, definition, features, point):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW.inches, SH.inches, WHITE)
    # 상단 밴드
    rect(s, 0, 0, SW.inches, 1.32, NAVY)
    rect(s, 0, 1.32, SW.inches, 0.05, MINT)
    textbox(s, 0.62, 0.26, 9.6, 0.7, [[(title, 26, WHITE, True, False)]])
    textbox(s, 0.64, 0.92, 9.6, 0.34, [[(route, 11.5, LTBLUE, False, False)]])
    # 역할 칩(아웃라인) — 관리자
    col = ADMIN if role == "관리자" else NAVY
    role_chip_outline(s, 11.55, 0.42, role + " 화면", col)
    # 페이지 번호(추후 renumber)
    textbox(s, 12.2, 6.95, 0.9, 0.4, [[("00 / 00", 10, GRAY, False, False)]], align=PP_ALIGN.RIGHT)

    # 좌측 스크린샷 플레이스홀더(빈 영역)
    ph = rect(s, IMG_L, IMG_T, IMG_W, IMG_H, LIGHTER, line=BORDER, round_=True)
    tf = ph.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "화면 스크린샷 영역"; _set(r, 13, FAINT, bold=True)

    # 우측 텍스트 — 정의
    rw = RIGHT_EDGE - RIGHT_X
    textbox(s, RIGHT_X, 1.6, 2, 0.34, [[("정의", 12, MINT, True, False)]])
    textbox(s, RIGHT_X, 1.96, rw, 0.9, [[(definition, 13, GRAY, False, False)]])
    # 핵심 기능
    textbox(s, RIGHT_X, 2.86, 4, 0.34, [[("핵심 기능", 12, BLUE, True, False)]])
    fy = 3.28
    for head, body in features:
        rect(s, RIGHT_X + 0.04, fy + 0.06, 0.12, 0.12, BLUE)
        textbox(s, RIGHT_X + 0.33, fy - 0.04, rw - 0.33, 0.72,
                [[(head + "  ", 13, INK, True, False), (body, 11, GRAY, False, False)]])
        fy += 0.72
    # POINT (하단 전체 폭)
    box = rect(s, 0.62, 6.05, 12.1, 0.82, LIGHT, line=RGBColor(0xBA, 0xD6, 0xF4), round_=True)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "POINT  "; _set(r, 12, MINT, bold=True)
    r = p.add_run(); r.text = point; _set(r, 13, NAVY, bold=False)
    return s


NEW_SLIDES = [
    ("VoC 애널리틱스 · 실시간 이슈 모니터링", "관리자", "/voc-console · monitor",
     "실시간 VoC 인입에서 급증 이슈·위험 키워드를 조기에 탐지하는 관제 화면입니다.",
     [("실시간 이슈 알림", "부지급 불만·'금감원' 키워드 급증 등 이상 징후를 카루셀 경보로 표시"),
      ("급증 키워드·유형", "키워드 TOP 10·클라우드 + 급증 VoC 유형(건수·전일 대비)"),
      ("부서별 처리 현황", "6개 부서 처리/유입·처리율과 병목·주의·정상 상태"),
      ("채널·기간 필터", "콜센터·이메일·대외기관 채널 및 기간별로 재집계")],
     "이상 징후를 실시간으로 포착해 민원이 확산되기 전에 선제 대응합니다."),

    ("VoC 애널리틱스 · 통계 대시보드", "관리자", "/voc-console · statsboard",
     "전사 VoC를 매일 전일 데이터로 집계하고, AI가 섹션별 인사이트를 제시하는 통계 화면입니다.",
     [("일일 현황", "문의·VoC·고위험 KPI(스파크라인) · VoC 구성 도넛 · 리스크 단계 퍼널"),
      ("처리 실적", "처리율·SLA·품질 게이지 + 부서별 처리 현황(오늘 vs 전일)"),
      ("기간 추세·유형 분석", "탐지 추이·인입 강도 히트맵 · 고객 프로파일·상품별 점유율"),
      ("AI 종합 분석", "섹션마다 위험 요인·병목·권고 사항을 자동 요약")],
     "흩어진 VoC 지표를 하나의 대시보드로 집계하고 AI가 해석까지 덧붙입니다."),

    ("VoC 애널리틱스 · VoC 리포트", "관리자", "/voc-console · report",
     "일일·월간 VoC 리포트를 자동 생성·시각화하고 배포하는 화면입니다.",
     [("리포트 설정", "일일/월간 · 기간·센터·채널 · 포함 섹션 선택"),
      ("자동 구성", "접수 추이·고객 경험 분포·주요 VoC·개선 사례를 문서형으로 편성"),
      ("금감원 민원 집계", "제기·처리 건수와 주요 민원 요약 포함"),
      ("내보내기", "PDF · 엑셀 · 인쇄 · 공유")],
     "반복 작성하던 VoC 리포트를 클릭 몇 번으로 표준 서식에 맞춰 생성합니다."),

    ("민원 탐지·이관", "관리자", "/complaint-detection",
     "콜·이메일·챗봇 문의를 단일 인입 큐로 통합해 유형 분류부터 부서 이관까지 처리하는 화면입니다.",
     [("통합 현황 대시보드", "VoC 비중·고위험 비중·이관 진행률·부서별 처리 현황"),
      ("통합 인입 큐", "미배정/배정완료/처리완료 탭 · 위험도·긴급도·담당 부서"),
      ("AI 분류·부서 이관", "유형 자동 분류·부서 추천 → 재배정·유형 변경"),
      ("VoC 등록·처리", "감지 키워드·트리거 확인 후 표준 형식으로 등록")],
     "흩어진 채널의 문의를 한 큐로 모아 분류·평가·부서 이관을 일괄 처리합니다."),

    ("대외민원 처리", "관리자", "/external-complaint",
     "금융감독원 등 대외기관 이첩 민원의 회신 초안을 작성·검증·처리하는 화면입니다.",
     [("금감원 접수함", "미처리·기한 임박(D-Day) 우선 정렬 · 처리 구분(이첩·자율조정·사실조회)"),
      ("원천 시스템 조회", "계약원장·청구심사·손해사정·CRM 사실관계 자동 집계"),
      ("자동 초안 작성", "금감원·민원인 회신 템플릿 + 검증 체크리스트(근거·금지표현·기한)"),
      ("근거·사례 첨부", "약관·법령·유사 분쟁조정 판례를 초안에 자동 인용")],
     "사실관계 조회부터 근거 인용·초안 작성까지, 대외기관 회신을 표준 절차로 지원합니다."),
]


def main():
    before = len(prs.slides._sldIdLst)
    for sl in NEW_SLIDES:
        new_slide(*sl)

    # 새 슬라이드를 기대 효과(마지막 2장) 앞으로 이동
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    new_ids = ids[before:]
    insert_at = before - 2  # 기대 효과 2장 앞
    for x in new_ids:
        sldIdLst.remove(x)
    for i, x in enumerate(new_ids):
        sldIdLst.insert(insert_at + i, x)

    # 페이지 번호 일괄 재부여 (NN / NN 패턴 텍스트박스)
    slides = list(prs.slides)
    total = len(slides)
    pat = re.compile(r"^\s*\d{1,2}\s*/\s*\d{1,2}\s*$")
    for pos, slide in enumerate(slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if pat.match(sh.text_frame.text.strip()):
                run = sh.text_frame.paragraphs[0].runs[0]
                run.text = f"{pos:02d} / {total:02d}"
                break

    prs.save(OUT)
    print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst),
          "| inserted:", len(NEW_SLIDES), "at index", insert_at)


if __name__ == "__main__":
    main()
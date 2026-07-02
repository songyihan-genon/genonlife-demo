# -*- coding: utf-8 -*-
"""데모 PPT 재배치: 스크린샷 좌측(크게) / 텍스트 우측 / POINT 하단 전체 / 엘보 커넥터(끝 동그라미). 원본 덮어쓰기."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

FILE = "GenON_AICC_데모_.pptx"

BLUE = RGBColor(0x2F, 0x6B, 0xB0)
BORDER = RGBColor(0xCD, 0xD9, 0xE8)

# 좌측 스크린샷 영역(인치) — 최대한 크게
IMG_L, IMG_T, IMG_W, IMG_H = 0.4, 1.55, 7.5, 4.5
# 우측 텍스트 영역
RIGHT_X = 8.1          # 본문 좌측 기준점(원래 0.62 → 이 값으로 이동)
RIGHT_EDGE = 12.9      # 우측 텍스트 한계
SHIFT = RIGHT_X - 0.62
BODY_TOP, BODY_BOT = 1.5, 6.0  # 우측으로 옮길 본문 영역(POINT는 6.05라 제외)

prs = Presentation(FILE)


def inch(emu):
    return Emu(emu or 0).inches


def add_circle_end(conn):
    ln = conn.line._get_or_add_ln()
    end = ln.makeelement(qn("a:tailEnd"), {"type": "oval", "w": "lg", "len": "lg"})
    ln.append(end)


changed = []
for si, slide in enumerate(prs.slides, 1):
    pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pics:
        continue
    pic = max(pics, key=lambda s: s.width * s.height)

    # 1) 스크린샷 좌측 + 영역에 맞춰 비율 유지 최대 크기
    w_in, h_in = inch(pic.width), inch(pic.height)
    scale = min(IMG_W / w_in, IMG_H / h_in)
    nw, nh = w_in * scale, h_in * scale
    pic.width, pic.height = Inches(nw), Inches(nh)
    pic.left = Inches(IMG_L + (IMG_W - nw) / 2)
    pic.top = Inches(IMG_T + (IMG_H - nh) / 2)
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(1)

    # 2) 본문 텍스트/마커를 우측으로 이동(+ 폭 재계산). POINT/페이지번호/상단밴드는 제외.
    bullets = []
    for sh in slide.shapes:
        if sh is pic:
            continue
        L, T, W, H = inch(sh.left), inch(sh.top), inch(sh.width), inch(sh.height)
        if not (BODY_TOP < T < BODY_BOT):
            continue
        new_left = L + SHIFT
        sh.left = Inches(new_left)
        if W < 0.25 and H < 0.25:   # 불릿 사각형
            bullets.append(inch(sh.top) + H / 2)
        else:                        # 텍스트 박스 → 우측 폭에 맞춤
            sh.width = Inches(max(1.5, RIGHT_EDGE - new_left))

    # 3) 핵심기능 불릿 → 스크린샷으로 엘보 커넥터(끝 동그라미). 끝점은 균등 분포(추후 드래그 미세조정)
    bullets.sort()
    n = len(bullets)
    for k, by in enumerate(bullets):
        ey = IMG_T + (k + 0.5) / max(n, 1) * IMG_H
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW,
            Inches(RIGHT_X), Inches(by),                 # 시작: 우측 텍스트 불릿
            Inches(IMG_L + IMG_W * 0.55), Inches(ey),    # 끝: 스크린샷 위(동그라미)
        )
        conn.line.color.rgb = BLUE
        conn.line.width = Pt(1.25)
        add_circle_end(conn)

    changed.append((si, n))

prs.save(FILE)
print("overwritten:", FILE)
for si, n in changed:
    print(f"  slide {si}: 스크린샷 좌측 정렬 + 커넥터 {n}개")